#!/usr/bin/env python3
"""
Assemble Figure S2 (v5) from its panels into PDF, PNG and PPTX -- de novo.

The MTG counterpart of fig_1/v3: figs2_v5.R is a copy of fig1_v3.R with the
inputs swapped region-for-region, and this is a copy of that figure's assembler
with the title and captions swapped to match. Layout is deliberately shared, so
Fig. S2 lands on the page exactly where Fig. 1 does.

Replaces the media-swap assemble_figure.py, which copied the hand-laid-out
fig_s2/v4.pptx and substituted its embedded images -- and pulled its coordinates
out of fig_1/v2/v2.pptx, making two superseded, untracked templates build
inputs. Neither is needed now: the placements below are the layout.

Two things that version needed and this one does not:
  * Panel D was cropped to its ink box with ghostscript and converted via
    pdf2svg. content_bbox() does that directly, so panel_D.pdf is used as-is.
  * SVG fill-opacity was pre-blended because PowerPoint's SVG importer drops
    semi-transparent fills. The PPTX now embeds rasters, so the blending has
    already happened at render time.

Layout (US Letter portrait, 612 x 792 pt):
  a  panel_A.svg              Jorstad markers
  b  panel_B.svg              Gabitto markers
  c  panel_C.svg              UpSet overlap  (nudged down to clear its caption)
  d  panel_D.pdf              unique/reproducible subclass markers table
  e/f panel_EF_combined.svg   metacell correlation panels (one image, two letters)
"""
import os, re, subprocess, tempfile
from collections import Counter

import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt

V5 = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.join(V5, "Kang_Figure_S2_v5")
OUT_PDF, OUT_PNG, OUT_PPTX = BASE + ".pdf", BASE + ".png", BASE + ".pptx"

# --- page geometry (points; US Letter portrait) -------------------------------
PW, PH = 612.0, 792.0
MARG = 7.1

TITLE = ("Fig. S2 | Unique subclass marker genes in human MTG vary by choice of "
         "study and algorithm")
TITLE_TOP, TITLE_H, TITLE_FS = 2.8, 24.0, 12.0
LET_FS, CAP_FS, CY_FS = 10.0, 6.0, 5.0

# PowerPoint insets its shape text; the coordinates below are shape origins from
# v2.pptx, so the same offsets fig_3/v4 needed apply here. See that script.
INSET_X, INSET_Y_BOX, INSET_Y_LETTER = 7.2, 5.0, 6.25

SRC = {
    "A":  "panel_A.svg",
    "B":  "panel_B.svg",
    "C":  "panel_C.svg",
    "D":  "panel_D.pdf",
    "EF": "panel_EF_combined.svg",
}

# Placement cells (x, y, w, h) -- the v2 picture frames, EMU/12700.
# Panel C carries the +10.8 pt (0.15 in) nudge the media-swap script applied so
# the plot clears its "Overlap of unique subclass markers" caption.
CELLS = {
    "A":  ( 49.0,  43.2, 244.8,  73.4),
    "B":  ( 49.0, 113.8, 244.8,  73.4),
    "C":  (312.5,  64.8, 218.1, 118.7),
    "D":  ( 45.4, 199.2, 150.5, 205.2),
    "EF": (216.4, 204.7, 392.8, 181.4),
}

LETTERS = {
    "a": ( 41.2,  29.1),
    "b": ( 41.0,  98.6),
    "c": (350.5,  44.4),
    "d": ( 19.8, 177.2),
    "e": (212.2, 176.0),
    "f": (401.8, 175.7),
}

# (text, x, y, w, centred?)
CAPTIONS = [
    ("Jorstad et al. 2023 (MTG, 3 donors)",  104.7,  36.0, 109.9, False),
    ("Gabitto et al. 2024 (MTG, 9 donors)",  104.5, 108.9, 109.6, False),
    ("Overlap of unique subclass markers",   399.5,  51.2, 126.0, True),
    ("Pairwise correlations of metacells \n(Jorstad et al. 2023, n = 30,284 genes)",
                                             233.4, 180.6, 117.4, True),
    ("Pairwise correlations of metacells \n(Gabitto et al. 2024, n = 36,601 genes)",
                                             419.8, 180.7, 117.0, True),
    ("Unique and reproducible subclass markers", 59.1, 185.4, 127.9, False),
]

# UpSetR's own y-axis title is suppressed (its placement is not controllable), so
# it is drawn here, rotated, in the empty space left of the bars. Positioned as a
# fraction of panel C's placed box so it follows the panel if that changes.
CY_TITLE = {"text": "# of intersecting genes", "fx": 0.329, "fy": 0.147}

# "n = <count> genes" in the E/F captions is rewritten from this sidecar, so the
# counts always match the data actually plotted (fig1_v3.R writes it).
GENE_COUNTS_FILE = "panel_EF_gene_counts.txt"


def read_gene_counts():
    d = {}
    try:
        for line in open(os.path.join(V5, GENE_COUNTS_FILE)):
            if "=" in line:
                k, v = line.strip().split("=", 1)
                d[k] = int(v)
    except FileNotFoundError:
        pass
    return d


def ef_body_centers(svg_path):
    """E and F heatmap body centres, as fractions of the SVG width.

    Located from the Set1 cell-class colour strip: collect its rects, split them
    at the widest horizontal gap, and take each group's span centre. Lets the f
    letter and caption track panel F instead of sitting at a fixed offset.
    """
    s = open(svg_path, encoding="utf-8", errors="ignore").read()
    VW = float(re.search(r"viewBox='0 0 ([\d.]+) ", s).group(1))
    S1 = ("#4DAF4A", "#377EB8", "#E41A1C")
    rects = [(float(m[1]), float(m[2]), float(m[3]))
             for m in re.finditer(r"<rect x='([\d.]+)' y='([\d.]+)' width='([\d.]+)' "
                                  r"height='[\d.]+' style='[^']*fill: (#[0-9A-Fa-f]{6})", s)
             if m[4].upper() in S1]
    strip_y = Counter(round(r[1], 1) for r in rects).most_common(1)[0][0]
    strip = sorted(r for r in rects if abs(r[1] - strip_y) < 0.5)
    gi = max(range(len(strip) - 1), key=lambda i: strip[i + 1][0] - (strip[i][0] + strip[i][2]))
    cen = lambda g: (min(r[0] for r in g) + max(r[0] + r[2] for r in g)) / 2 / VW
    return cen(strip[:gi + 1]), cen(strip[gi + 1:])


def stage(tmp):
    """Panel sources -> PDF. Inkscape rather than rsvg (see fig_3/v4)."""
    out = {}
    for k, name in SRC.items():
        src = os.path.join(V5, name)
        if name.endswith(".pdf"):
            out[k] = src
            continue
        p = os.path.join(tmp, f"{k}.pdf")
        subprocess.run(["inkscape", src, "--export-type=pdf",
                        f"--export-filename={p}"], check=True, capture_output=True)
        out[k] = p
    return out


def content_bbox(src_pdf):
    page = fitz.open(src_pdf)[0]
    r = fitz.Rect(1e9, 1e9, -1e9, -1e9)
    for b in page.get_text("blocks"):
        r |= fitz.Rect(b[:4])
    for d in page.get_drawings():
        r |= d["rect"]

    # get_drawings() does not descend into Form XObjects, and Inkscape puts content
    # inside them -- fig_4/v6's panel S had its whole column dendrogram clipped away
    # by trimming to the under-reported box. Union with a raster ink scan, which sees
    # whatever the page actually paints; the union can only grow the box, never crop.
    # Applied as a SAFETY NET, not a co-equal source: the scan rounds outward by up
    # to a pixel and antialiasing bleeds past the true vector edge, so taking it
    # verbatim would nudge every box by a fraction of a point and shift figures that
    # are already published. Only genuinely missed content -- beyond SCAN_TOL -- wins.
    SCAN_DPI, SCAN_TOL = 144, 1.0
    pm = page.get_pixmap(dpi=SCAN_DPI)
    a = np.frombuffer(pm.samples, dtype=np.uint8).reshape(pm.height, pm.width, pm.n)[:, :, :3]
    ys, xs = np.where((a < 250).any(2))
    if len(ys):
        k = 72.0 / SCAN_DPI
        ink = fitz.Rect(page.rect.x0 + xs.min() * k, page.rect.y0 + ys.min() * k,
                        page.rect.x0 + (xs.max() + 1) * k, page.rect.y0 + (ys.max() + 1) * k)
        if ink.x0 < r.x0 - SCAN_TOL: r.x0 = ink.x0
        if ink.y0 < r.y0 - SCAN_TOL: r.y0 = ink.y0
        if ink.x1 > r.x1 + SCAN_TOL: r.x1 = ink.x1
        if ink.y1 > r.y1 + SCAN_TOL: r.y1 = ink.y1

    if r.is_empty or r.x1 < r.x0:
        return page.rect
    return r & page.rect


def fit(bb, cx, cy, cw, ch):
    """Aspect-preserving fit, centred in the cell. The media-swap script
    preserved each picture's top-left corner, but centring reproduces the
    rendered result more closely (2.3% vs 4.2% blurred pixel difference against
    fig1_v3_assembled.pdf), so the template's boxes evidently already allow for
    the centring PowerPoint applies."""
    s = min(cw / bb.width, ch / bb.height)
    dw, dh = bb.width * s, bb.height * s
    return fitz.Rect(cx + (cw - dw) / 2, cy + (ch - dh) / 2,
                     cx + (cw - dw) / 2 + dw, cy + (ch - dh) / 2 + dh)


def compute_layout(pdf):
    panels = []
    for k, cell in CELLS.items():
        bb = content_bbox(pdf[k])
        panels.append((k, pdf[k], bb, fit(bb, *cell)))
    cbox = next(r for k, _, _, r in panels if k == "C")
    cy = (cbox.x0 + CY_TITLE["fx"] * cbox.width,
          cbox.y0 + CY_TITLE["fy"] * cbox.height)

    letters = dict(LETTERS)
    captions = [list(c) for c in CAPTIONS]

    # Region-correct gene counts from the sidecar.
    counts = read_gene_counts()
    for c in captions:
        for ds, key in (("Jorstad", "jorstad"), ("Gabitto", "gabitto")):
            if ds in c[0] and "n = " in c[0] and key in counts:
                c[0] = re.sub(r"n\s*=\s*[\d,]+\s*genes",
                              f"n = {counts[key]:,} genes", c[0])

    # f letter and caption track panel F's body rather than sitting at a fixed
    # offset, so they stay aligned if the E/F panel is regenerated.
    efbox = next(r for k, _, _, r in panels if k == "EF")
    try:
        eF, fF = ef_body_centers(os.path.join(V5, SRC["EF"]))
        e_cen, f_cen = efbox.x0 + eF * efbox.width, efbox.x0 + fF * efbox.width
        letters["f"] = (letters["e"][0] + (f_cen - e_cen), letters["f"][1])
        eT = next(c for c in captions if "Pairwise" in c[0] and "Jorstad" in c[0])
        fT = next(c for c in captions if "Pairwise" in c[0] and "Gabitto" in c[0])
        fT[1] = eT[1] + (f_cen - e_cen)
    except Exception as exc:                      # strip not found -> keep template x
        print(f"  note: E/F body centres unavailable ({exc}); using template x")

    return {"panels": panels, "cy_centre": cy,
            "letters": letters, "captions": captions}


def render_pdf(lay):
    doc = fitz.open()
    page = doc.new_page(width=PW, height=PH)
    page.insert_textbox(fitz.Rect(MARG + INSET_X, TITLE_TOP + INSET_Y_BOX,
                                  PW - MARG, TITLE_TOP + INSET_Y_BOX + TITLE_H),
                        TITLE, fontname="hebo", fontsize=TITLE_FS,
                        align=fitz.TEXT_ALIGN_LEFT, lineheight=1.15)
    for _, src, bb, rect in lay["panels"]:
        page.show_pdf_page(rect, fitz.open(src), 0, clip=bb)
    for ch, (x, y) in lay["letters"].items():
        page.insert_text((x + INSET_X, y + INSET_Y_LETTER + LET_FS), ch,
                         fontname="hebo", fontsize=LET_FS)
    for text, x, y, w, centred in lay["captions"]:
        x0 = x if centred else x + INSET_X          # centred text centres in the raw box
        page.insert_textbox(
            fitz.Rect(x0, y + INSET_Y_BOX, x0 + w, y + INSET_Y_BOX + 4 * CAP_FS), text,
            fontname="helv", fontsize=CAP_FS, lineheight=1.2,
            align=fitz.TEXT_ALIGN_CENTER if centred else fitz.TEXT_ALIGN_LEFT)
    cx, cy = lay["cy_centre"]
    tw = fitz.get_text_length(CY_TITLE["text"], "helv", CY_FS)
    page.insert_textbox(fitz.Rect(cx - CY_FS, cy - tw / 2 - 2,
                                  cx + CY_FS * 1.6, cy + tw / 2 + 2),
                        CY_TITLE["text"], fontname="helv", fontsize=CY_FS,
                        rotate=90, align=fitz.TEXT_ALIGN_CENTER)
    doc.save(OUT_PDF, deflate=True)
    return doc


def render_png(doc):
    doc[0].get_pixmap(dpi=300).save(OUT_PNG)


def _box(slide, x, y, w, h, rot=None):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    if rot is not None:
        tb.rotation = rot
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, bold=False):
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    return r


def render_pptx(lay, tmp):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(PW), Pt(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for key, src, bb, rect in lay["panels"]:
        png = os.path.join(tmp, f"{key}_raster.png")
        fitz.open(src)[0].get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72),
                                     clip=bb).save(png)
        slide.shapes.add_picture(png, Pt(rect.x0), Pt(rect.y0),
                                 Pt(rect.width), Pt(rect.height))

    _run(_box(slide, MARG + INSET_X, TITLE_TOP + INSET_Y_BOX,
              PW - 2 * MARG, TITLE_H).paragraphs[0], TITLE, TITLE_FS, bold=True)
    for ch, (x, y) in lay["letters"].items():
        _run(_box(slide, x + INSET_X, y + INSET_Y_BOX, 20.0,
                  LET_FS * 1.7).paragraphs[0], ch, LET_FS, bold=True)
    for text, x, y, w, centred in lay["captions"]:
        _run(_box(slide, x if centred else x + INSET_X, y + INSET_Y_BOX, w,
                  4 * CAP_FS).paragraphs[0], text, CAP_FS)
    cx, cy = lay["cy_centre"]
    _run(_box(slide, cx - 25, cy - 4, 50, 8, rot=270).paragraphs[0],
         CY_TITLE["text"], CY_FS)

    prs.save(OUT_PPTX)


def main():
    with tempfile.TemporaryDirectory() as tmp:
        pdf = stage(tmp)
        lay = compute_layout(pdf)
        doc = render_pdf(lay)
        render_png(doc)
        render_pptx(lay, tmp)
    for p in (OUT_PDF, OUT_PNG, OUT_PPTX):
        print(f"wrote {p}  ({os.path.getsize(p)/1e6:.2f} MB)")


if __name__ == "__main__":
    main()
