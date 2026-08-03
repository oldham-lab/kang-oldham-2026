#!/usr/bin/env python
# Assemble Figure 5 (v5) into a PowerPoint slide from four panels:
#   a = panel_A.svg  (REI schematic, provided)
#   b = panel_B.svg  (labeled dendrogram + barplots)
#   c = panel_C.svg  (Jorstad heatmap, row dendrogram, no x-axis)
#   d = panel_D.svg  (Gabitto heatmap, row dendrogram, x-axis)
#   + heatmap colour legend (panel_C_consensusMin_Jorstad_legend.svg)
#
# Column alignment: each heatmap/dendrogram writes its column band (x0,x1) in R.
# b/c/d are placed so their band maps to the SAME slide x-range [BAND_L, BAND_R], so the
# meta-modules line up vertically even though the panels have different left/right margins
# (panel B's band is narrower because of the "% of all mods/genes" annotation text).
#
# Rasterisation: Inkscape (NOT rsvg-convert), then flatten alpha onto white -- rsvg renders
# panel A's coloured schematic as solid black boxes; Inkscape renders it faithfully (per
# fig_3/v4/assemble_figure.py). Run with a python that has python-pptx (e.g. the venv).
import os
import subprocess, os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

V5  = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_5/v5")
DPI = 300
OUT = f"{V5}/Kang_Figure_5_v5.pptx"

def _sh(cmd): subprocess.run(cmd, shell=True, check=True, capture_output=True)

def svg2png(stem, src=None):
    src = src or f"{V5}/{stem}.svg"; png = f"{V5}/{stem}.png"
    _sh(f"inkscape {src} --export-type=png --export-filename={png} "
        f"--export-dpi={DPI} --export-background=white --export-background-opacity=1")
    _sh(f"convert {png} -background white -alpha remove -alpha off {png}")  # flatten alpha
    return png

def aspect(png):                       # height / width
    with Image.open(png) as im: return im.height / im.width

def band_frac(coords):                 # (f0, f1) column-band fractions of image width
    d = {k: float(v) for k, v in (l.split("=") for l in open(coords))}
    return d["x0"] / d["W_px"], d["x1"] / d["W_px"]

# ---- rasterise panels + read column bands ----
pA = svg2png("panel_A"); pB = svg2png("panel_B"); pC = svg2png("panel_C"); pD = svg2png("panel_D")
pL = svg2png("legend", src=f"{V5}/panel_C_consensusMin_Jorstad_legend.svg")
fB = band_frac(f"{V5}/jorstad_panel_coords.txt")
fC = band_frac(f"{V5}/panel_C_coords.txt")
fD = band_frac(f"{V5}/panel_D_coords.txt")

# ---- layout constants (inches), matched to v4/Kang_Figure_5_v4.pptx ----
BAND_L, BAND_R = 1.23, 7.61            # v4 heatmap column band on the page
B_TOP, D_TOP = 2.55, 7.80             # tops of panels b and d (match v4 spacing)
# C_TOP is derived so the b|c gap equals the c|d gap (see below).
# v4 panel-letter, subtitle, panel-A and legend geometry (EMU->in from the v4 pptx)
A_BOX  = (2.72, 0.484, 3.47, 1.75)     # panel A picture box L,T,W,H
LEG    = (7.67, 6.744, 0.75, 1.99)     # legend box L,T,W,H

def band_h(png, frac):                 # placed height when band -> [BAND_L, BAND_R]
    f0, f1 = frac
    return (BAND_R - BAND_L) / (f1 - f0) * aspect(png)

def place_by_band(png, frac, top):
    """Scale/position so the panel's column band spans [BAND_L, BAND_R]; return bottom."""
    f0, f1 = frac
    W = (BAND_R - BAND_L) / (f1 - f0)
    L = BAND_L - f0 * W
    slide.shapes.add_picture(png, Inches(L), Inches(top), Inches(W), Inches(W * aspect(png)))
    return top + W * aspect(png)

# ---- slide ----
prs = Presentation(); prs.slide_width = Inches(8.5); prs.slide_height = Inches(11.0)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def textbox(x, y, w, h, text, size, bold=True, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = "Arial"; r.font.color.rgb = RGBColor(0,0,0)
    return tb

def subtitle(text, x, y):               # centered over the column band, anchored near v4 y
    textbox(BAND_L, y, BAND_R - BAND_L, 0.30, text, 8, bold=False, align=PP_ALIGN.CENTER)

# --- pictures first (so text labels placed afterwards sit on top) ---
aL, aT, aW, aH = A_BOX
slide.shapes.add_picture(pA, Inches(aL), Inches(aT), Inches(aW), Inches(aW * aspect(pA)))
b_bot = place_by_band(pB, fB, B_TOP)
# shift panel c down so the b|c gap equals the c|d gap
C_TOP = (b_bot + D_TOP - band_h(pC, fC)) / 2
c_bot = place_by_band(pC, fC, C_TOP)
d_bot = place_by_band(pD, fD, D_TOP)
lL, lT, lW, lH = LEG
slide.shapes.add_picture(pL, Inches(lL), Inches(lT), height=Inches(lH))   # legend, straddling c/d

# --- text on top ---
textbox(0.1, 0.035, 8.2, 0.45,
        "Fig. 5 | Clustering CoPA projection patterns reveals the major motifs of gene "
        "activity in human dorsal frontal cortex (DFC) cell types", 10)
textbox(2.578, 0.35, 0.35, 0.33, "a", 14)
textbox(0.767, 2.279, 0.35, 0.33, "b", 14)
subtitle("Consensus clustering of all modules over all subclasses "
         "(Jorstad et al. 2023 projections, DFC)", 1.931, 2.343)
textbox(0.767, C_TOP - 0.20, 0.35, 0.33, "c", 14)
textbox(0.77, 7.556, 0.35, 0.33, "d", 14)
subtitle("Gabitto et al. 2024 projections arranged using Jorstad et al. 2023 clusters", 2.398, 7.618)

prs.save(OUT); print(f"Saved {OUT}")
print(f"  band -> slide x [{BAND_L},{BAND_R}] in;  bottoms b={b_bot:.2f} c={c_bot:.2f} d={d_bot:.2f}")

# ---- PDF + PNG (LibreOffice -> PDF -> PNG); PNG panels embed fine (no SVG in pptx) ----
stem = os.path.splitext(OUT)[0]
_sh(f"soffice --headless -env:UserInstallation=file:///tmp/soffice_fig5_v5 "
    f"--convert-to pdf --outdir {V5} {OUT}")
_sh(f"pdftoppm -png -r {DPI} -singlefile {stem}.pdf {stem}")
print(f"Saved {stem}.pdf and {stem}.png")
