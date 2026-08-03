#!/usr/bin/env python
# Assemble Figure S8 (v2) -- the MTG analogue of Figure 5. Same script as fig_5/v5's
# assembler EXCEPT there is no panel A schematic (so the three panels are labeled a/b/c)
# and the data/region is MTG. Panels (files kept as B/C/D to mirror fig_5):
#   a = panel_B.svg  (labeled dendrogram + barplots)
#   b = panel_C.svg  (Jorstad MTG heatmap, row dendrogram, no x-axis)
#   c = panel_D.svg  (Gabitto MTG heatmap, row dendrogram, x-axis)
#   + heatmap colour legend
# b/c/a placed by their column band -> common slide x-range so meta-modules stay aligned;
# the middle panel's top is derived so the a|b gap equals the b|c gap. Geometry matched to
# v2/Kang_Figure_S8_v2.pptx. Rasterise with Inkscape (+flatten alpha) as in fig_3/fig_5.
# Run with a python that has python-pptx (e.g. /home/gugene/.venv/bin/python).
import os
import subprocess, os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

V   = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_s8/v2")
DPI = 300
OUT = f"{V}/Kang_Figure_S8_v2.pptx"

def _sh(cmd): subprocess.run(cmd, shell=True, check=True, capture_output=True)
def svg2png(stem, src=None):
    src = src or f"{V}/{stem}.svg"; png = f"{V}/{stem}.png"
    _sh(f"inkscape {src} --export-type=png --export-filename={png} "
        f"--export-dpi={DPI} --export-background=white --export-background-opacity=1")
    _sh(f"convert {png} -background white -alpha remove -alpha off {png}")
    return png
def aspect(png):
    with Image.open(png) as im: return im.height / im.width
def band_frac(coords):
    d = {k: float(v) for k, v in (l.split("=") for l in open(coords))}
    return d["x0"] / d["W_px"], d["x1"] / d["W_px"]

# panels (a=dendrogram+barplots, b=Jorstad heatmap, c=Gabitto heatmap) + legend
pa = svg2png("panel_B"); pb = svg2png("panel_C"); pc = svg2png("panel_D")
pL = svg2png("legend", src=f"{V}/panel_C_consensusMin_Jorstad_legend.svg")
fa = band_frac(f"{V}/jorstad_panel_coords.txt")
fb = band_frac(f"{V}/panel_C_coords.txt")
fc = band_frac(f"{V}/panel_D_coords.txt")

# ---- layout constants (inches), matched to v2/Kang_Figure_S8_v2.pptx ----
BAND_L, BAND_R = 1.33, 7.58            # MTG heatmap column band on the page
A_TOP, C_TOP   = 0.74, 5.95            # tops of panels a and c; b derived for equal gaps
LEG            = (7.66, 4.96, 0.75, 2.0)  # legend box L,T,W,H

def band_h(png, frac):
    f0, f1 = frac
    return (BAND_R - BAND_L) / (f1 - f0) * aspect(png)

prs = Presentation(); prs.slide_width = Inches(8.5); prs.slide_height = Inches(11.0)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def place(png, frac, top):
    f0, f1 = frac
    W = (BAND_R - BAND_L) / (f1 - f0)
    slide.shapes.add_picture(png, Inches(BAND_L - f0 * W), Inches(top), Inches(W), Inches(W * aspect(png)))
    return top + W * aspect(png)

def textbox(x, y, w, h, text, size, bold=True, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = "Arial"; r.font.color.rgb = RGBColor(0,0,0)
def subtitle(text, y):
    textbox(BAND_L, y, BAND_R - BAND_L, 0.30, text, 8, bold=False, align=PP_ALIGN.CENTER)

# pictures first (text on top afterwards)
a_bot = place(pa, fa, A_TOP)
B_TOP = (a_bot + C_TOP - band_h(pb, fb)) / 2     # equal a|b and b|c gaps
b_bot = place(pb, fb, B_TOP)
c_bot = place(pc, fc, C_TOP)
lL, lT, lW, lH = LEG
slide.shapes.add_picture(pL, Inches(lL), Inches(lT), height=Inches(lH))

# text
textbox(0.1, 0.035, 8.0, 0.5,
        "Fig. S8 | Clustering CoPA projection patterns reveals the major motifs of gene "
        "activity in human medial temporal gyrus (MTG) cell types", 10)
textbox(0.61, A_TOP - 0.24, 0.35, 0.33, "a", 14)
subtitle("Consensus clustering of all modules over all subclasses "
         "(Jorstad et al. 2023 projections, MTG)", A_TOP - 0.22)
textbox(0.61, B_TOP - 0.20, 0.35, 0.33, "b", 14)
textbox(0.61, C_TOP - 0.24, 0.35, 0.33, "c", 14)
subtitle("Gabitto et al. 2024 projections arranged using Jorstad et al. 2023 MTG clusters", C_TOP - 0.22)

prs.save(OUT); print(f"Saved {OUT}")
print(f"  band x [{BAND_L},{BAND_R}]; B_TOP={B_TOP:.2f}; bottoms a={a_bot:.2f} b={b_bot:.2f} c={c_bot:.2f}")

# ---- PDF + PNG ----
stem = os.path.splitext(OUT)[0]
_sh(f"soffice --headless -env:UserInstallation=file:///tmp/soffice_s8_v2 --convert-to pdf --outdir {V} {OUT}")
_sh(f"pdftoppm -png -r {DPI} -singlefile {stem}.pdf {stem}")
print(f"Saved {stem}.pdf and {stem}.png")
