#!/usr/bin/env python3
"""
Figure assembly for Fig. 8 (v2).

Deliverables (all named Kang_Figure_8_v2.*):
  .pdf   — vector
  .png   — raster, 300 dpi
  .pptx  — one slide, the figure embedded as a vector SVG with a PNG fallback

Layout reference = the hand-laid-out fig_8/v2/v2.pptx (panel positions/sizes, the
panel letters a–f, the per-panel titles). This script reads those positions, then
composes a single master SVG that INLINES the panel sources (so they stay vector)
with the letters / titles drawn as Arial <text>. The master SVG is an internal
intermediate (deleted at the end, not a deliverable): rsvg-convert renders it to
the PDF and PNG, and it is embedded as the pptx's vector layer.

Panel sources live in fig_8/v1 (single source of truth). The template's embedded
copies are ignored; v1 is re-inlined so the figure tracks the current v1 files:
  panel a (top-left scatter)   <- panel_A_indiv_DFC.svg
  panel b (top-right scatter)  <- panel_B_indiv_DFC.svg
  panel c (overlap heatmap)    <- panel_C_D_DFC_panel_lower_nolegend.svg
  panel c legend               <- panel_C_D_DFC_legend.svg
  panel d (cell-type overlap)  <- panel_E_F_DFC_lower.svg
  panel e (GSEA summary table) <- panel_G_gsea_summary_table_dfc_fdr.pdf  (PDF ->
                                  cropped vector SVG via Inkscape)
  panel f (bottom-right)       <- LEFT BLANK for now (only the 'f' letter is drawn)
"""
import os, re, subprocess
import xml.etree.ElementTree as ET
import zipfile

V2       = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_8/v2")
V1       = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_8/v1")
TEMPLATE = f"{V2}/v2.pptx"
OUT_PDF  = f"{V2}/Kang_Figure_8_v2.pdf"
OUT_PNG  = f"{V2}/Kang_Figure_8_v2.png"
OUT_PPTX = f"{V2}/Kang_Figure_8_v2.pptx"
# internal intermediates (deleted at the end; not deliverables)
MASTER_SVG  = f"{V2}/.fig8_master.svg"
PANEL_E_SVG = f"{V2}/.fig8_panel_e.svg"

# template media basename -> (kind, source path). 'svg' is inlined directly;
# 'pdf' is converted to a cropped vector SVG first. image14 (panel f) is absent
# on purpose -> that panel is left blank.
MEDIA_SRC = {
    "image10.svg": ("svg", f"{V1}/panel_A_indiv_DFC.svg"),                 # a
    "image8.svg":  ("svg", f"{V1}/panel_B_indiv_DFC.svg"),                 # b
    "image4.svg":  ("svg", f"{V1}/panel_C_D_DFC_panel_lower_nolegend.svg"),# c
    "image6.svg":  ("svg", f"{V1}/panel_C_D_DFC_legend.svg"),             # c legend
    "image2.svg":  ("svg", f"{V1}/panel_E_F_DFC_lower.svg"),               # d
    "image12.svg": ("pdf", f"{V1}/panel_G_gsea_summary_table_dfc_fdr.pdf"),# e
}
# render order (panels, then the small c-legend on top of panel c)
RENDER_ORDER = ["image10.svg", "image8.svg", "image4.svg", "image2.svg",
                "image12.svg", "image6.svg"]

EMU_PER_PX = 9525.0          # 96 dpi
def px(emu): return float(emu) / EMU_PER_PX
def ptpx(centipt): return float(centipt) / 100.0 * 96.0 / 72.0  # PowerPoint sz -> px

NS = {"a":"http://schemas.openxmlformats.org/drawingml/2006/main",
      "p":"http://schemas.openxmlformats.org/presentationml/2006/main",
      "r":"http://schemas.openxmlformats.org/officeDocument/2006/relationships",
      "asvg":"http://schemas.microsoft.com/office/drawing/2016/SVG/main"}

# ---- read layout from v2.pptx --------------------------------------------
z = zipfile.ZipFile(TEMPLATE)
rels = {rel.get("Id"): rel.get("Target").split("/")[-1]
        for rel in ET.fromstring(z.read("ppt/slides/_rels/slide1.xml.rels"))}
slide = ET.fromstring(z.read("ppt/slides/slide1.xml"))
pres  = ET.fromstring(z.read("ppt/presentation.xml"))

sldSz = pres.find("p:sldSz", NS)
W_EMU, H_EMU = int(sldSz.get("cx")), int(sldSz.get("cy"))   # full slide canvas
W, H = px(W_EMU), px(H_EMU)

pics = {}   # media basename -> (x, y, w, h) in EMU
for pic in slide.iter("{%s}pic" % NS["p"]):
    svg = pic.find(".//asvg:svgBlip", NS)
    if svg is None:
        continue
    media = rels.get(svg.get("{%s}embed" % NS["r"]))
    xf = pic.find(".//a:xfrm", NS); off = xf.find("a:off", NS); ext = xf.find("a:ext", NS)
    pics.setdefault(media, (int(off.get("x")), int(off.get("y")),
                            int(ext.get("cx")), int(ext.get("cy"))))

# text shapes: keep explicit paragraph breaks (the wrapping titles are split into
# one <a:p> per line in the template), plus size / bold / alignment.
texts = []  # (lines, x, y, cx, cy, size_px, bold, centred)
for sp in slide.iter("{%s}sp" % NS["p"]):
    lines = []
    for p_el in sp.findall(".//a:p", NS):
        line = "".join(t.text or "" for t in p_el.findall(".//a:t", NS)).strip()
        if line:
            lines.append(line)
    if not lines:
        continue
    xf = sp.find(".//a:xfrm", NS); off = xf.find("a:off", NS); ext = xf.find("a:ext", NS)
    rpr = sp.find(".//a:rPr", NS)
    sz  = ptpx(rpr.get("sz")) if (rpr is not None and rpr.get("sz")) else 12.0
    bold = bool(rpr is not None and rpr.get("b") == "1")
    centred = any(p.get("algn") == "ctr" for p in sp.findall(".//a:pPr", NS))
    texts.append((lines, int(off.get("x")), int(off.get("y")),
                  int(ext.get("cx")), int(ext.get("cy")), sz, bold, centred))
z.close()

# ---- panel e: PDF -> content-cropped vector SVG --------------------------
subprocess.run(["inkscape", MEDIA_SRC["image12.svg"][1],
                "--export-type=svg", "--export-area-drawing",
                f"--export-filename={PANEL_E_SVG}"],
               check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
SRC_PATH = {m: (PANEL_E_SVG if m == "image12.svg" else MEDIA_SRC[m][1])
            for m in MEDIA_SRC}

# ---- SVG inlining (keeps panels vector; prefix ids to avoid cross-panel clashes)
def _idprefix(s, p):
    s = re.sub(r'\bid="([^"]+)"',          lambda m: f'id="{p}{m.group(1)}"', s)
    s = re.sub(r"\bid='([^']+)'",          lambda m: f"id='{p}{m.group(1)}'", s)
    s = re.sub(r'xlink:href="#([^"]+)"',   lambda m: f'xlink:href="#{p}{m.group(1)}"', s)
    s = re.sub(r"xlink:href='#([^']+)'",   lambda m: f"xlink:href='#{p}{m.group(1)}'", s)
    s = re.sub(r'(?<!xlink:)\bhref="#([^"]+)"', lambda m: f'href="#{p}{m.group(1)}"', s)
    s = re.sub(r'url\(#([^)]+)\)',         lambda m: f'url(#{p}{m.group(1)})', s)
    return s

def inline(path, prefix, x, y, w, h):
    s = open(path, encoding="utf-8").read()
    s = re.sub(r'<\?xml[^>]*\?>', '', s)
    s = re.sub(r'<!DOCTYPE[^>]*>', '', s).strip()
    head = re.search(r'<svg\b([^>]*)>', s).group(1)
    vb = re.search(r'viewBox=["\']([^"\']+)["\']', head).group(1)
    # keep the source's xmlns:* declarations (Inkscape SVGs use rdf/sodipodi/
    # inkscape prefixes inside; dropping their decls breaks the XML parse).
    xmlns = " ".join(re.findall(r'xmlns:[\w-]+=["\'][^"\']*["\']', head))
    s = _idprefix(s, prefix)
    s = re.sub(r'<svg\b[^>]*>',
               f'<svg {xmlns} x="{x:.2f}" y="{y:.2f}" width="{w:.2f}" height="{h:.2f}" '
               f'viewBox="{vb}" preserveAspectRatio="xMidYMid meet" overflow="visible">',
               s, count=1)
    return s

def esc(t): return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
FAMILY = "Arial, 'Liberation Sans', sans-serif"

parts = [f'<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" '
         f'width="{W:.2f}" height="{H:.2f}" viewBox="0 0 {W:.2f} {H:.2f}">',
         f'<rect width="{W:.2f}" height="{H:.2f}" fill="white"/>']

# panels (vector), in render order
for media in RENDER_ORDER:
    x, y, w, h = pics[media]
    parts.append(inline(SRC_PATH[media], media.replace(".", "_") + "_",
                         px(x), px(y), px(w), px(h)))

# letters (a–f) and per-panel titles, drawn as Arial. Explicit paragraph lines are
# stacked; centred shapes are centred on their text-box, left ones flush-left.
for lines, x, y, cx, cy, sz, bold, centred in texts:
    weight = ' font-weight="bold"' if bold else ''
    if centred:
        tx, anchor = px(x + cx / 2), "middle"
    else:
        tx, anchor = px(x) + 2, "start"
    for i, line in enumerate(lines):
        ty = px(y) + sz * 0.95 + i * sz * 1.2     # box-top -> baseline, then line-height
        parts.append(f'<text x="{tx:.1f}" y="{ty:.1f}" text-anchor="{anchor}" '
                     f'font-family="{FAMILY}" font-size="{sz:.1f}"{weight}>{esc(line)}</text>')

parts.append("</svg>")
open(MASTER_SVG, "w", encoding="utf-8").write("\n".join(parts))   # intermediate only

# vector PDF (deliverable)
subprocess.run(["rsvg-convert", "-f", "pdf", MASTER_SVG, "-o", OUT_PDF], check=True)
print("wrote:", OUT_PDF, f"({W:.0f}x{H:.0f}px)")

# raster PNG (deliverable). The master SVG is sized in px at 96 dpi, so rsvg's
# --dpi is a no-op; scale the pixel dimensions to render at 300 dpi instead.
PNG_DPI = 300
png_w, png_h = round(W * PNG_DPI / 96), round(H * PNG_DPI / 96)
subprocess.run(["rsvg-convert", "-f", "png", "-w", str(png_w), "-h", str(png_h),
                MASTER_SVG, "-o", OUT_PNG], check=True)
print("wrote:", OUT_PNG, f"({png_w}x{png_h}px @ {PNG_DPI}dpi)")

# ---- PPTX (deliverable): one slide = the composed figure, embedded as a vector
# SVG with a PNG fallback (the native PowerPoint pattern: a:blip -> PNG,
# asvg:svgBlip -> SVG), so the slide stays vector in PowerPoint while remaining
# openable everywhere. The SVG lives only inside the pptx; no .svg file is emitted.
from pptx import Presentation
from pptx.util import Emu
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(int(W_EMU)), Emu(int(H_EMU))
slide = prs.slides.add_slide(prs.slide_layouts[6])          # blank layout
pic = slide.shapes.add_picture(OUT_PNG, 0, 0, Emu(int(W_EMU)), Emu(int(H_EMU)))

# add the master SVG as a media part related from the slide, then point an
# asvg:svgBlip at it inside the picture's a:blip (alongside the PNG fallback).
spart = slide.part
svg_partname = spart.package.next_partname("/ppt/media/image%d.svg")
svg_part = Part(svg_partname, "image/svg+xml",
                spart.package, open(MASTER_SVG, "rb").read())
IMAGE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
svg_rid = spart.relate_to(svg_part, IMAGE_REL)

blip = pic._element.blipFill.find(qn("a:blip"))
blip.append(parse_xml(
    '<a:extLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">'
    '<asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    f'r:embed="{svg_rid}"/></a:ext></a:extLst>'))

prs.save(OUT_PPTX)
for tmp in (MASTER_SVG, PANEL_E_SVG):   # intermediates; not deliverables
    os.remove(tmp)
print("wrote:", OUT_PPTX)
