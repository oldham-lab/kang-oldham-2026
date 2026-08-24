#!/usr/bin/env python3
"""
Figure assembly for Fig. 8 (v3).

Deliverables (all named Kang_Figure_8_v4.*):
  .pdf   — vector
  .png   — raster, 300 dpi
  .pptx  — one slide, the figure embedded as a vector SVG with a PNG fallback

Layout: the pics/texts tables below ARE the layout. They were read once out of the
hand-laid-out v2.pptx that used to sit in this folder and frozen here, so the
untracked template is no longer a build input and the figure rebuilds from a
clone. This script composes a single master SVG that INLINES the panel sources (so they stay vector)
with the letters / titles drawn as Arial <text>. The master SVG is an internal
intermediate (deleted at the end, not a deliverable): rsvg-convert renders it to
the PDF and PNG, and it is embedded as the pptx's vector layer.

Panel sources all live in this folder (MEDIA_SRC below); the keys are the template's
old media names, kept because the render order and the a/b special-casing refer to
them:
  panel a (top-left scatter)   <- panel_A_indiv_DFC.svg
  panel b (top-right scatter)  <- panel_B_indiv_DFC.svg
  panel c (overlap heatmap)    <- panel_C_D_DFC_panel_lower_nolegend.svg
  panel c legend               <- panel_C_D_DFC_legend.svg
  panel d (cell-type overlap)  <- panel_E_F_DFC_lower.svg
  panel e (GSEA summary table) <- v3/panel_e_gsea_table.pdf  (red highlight + ref
                                  superscripts removed; make_panel_e.py)
  panel f (module 992 snapshot)<- v3/992_DFC_Schizophrenia_control_bulk_megaset.svg
                                  (no red highlight; fig8_specific_modules.R)
PDF panel sources are converted to cropped vector SVGs via Inkscape before inlining.
"""
import os, re, subprocess

V4       = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_8/v4")
OUT_PDF  = f"{V4}/Kang_Figure_8_v4.pdf"
OUT_PNG  = f"{V4}/Kang_Figure_8_v4.png"
OUT_PPTX = f"{V4}/Kang_Figure_8_v4.pptx"
# internal intermediates (deleted at the end; not deliverables)
MASTER_SVG  = f"{V4}/.fig8_master.svg"
PANEL_E_SVG = f"{V4}/.fig8_panel_e.svg"

# template media basename -> (kind, source path). 'svg' is inlined directly;
# 'pdf' is converted to a cropped vector SVG first.
MEDIA_SRC = {
    "image10.svg": ("svg", f"{V4}/panel_A_indiv_DFC_nolegend.svg"),        # a (nolegend + fig_7-style r bracket; fig8_panelA_bracket.R)
    "image8.svg":  ("svg", f"{V4}/panel_B_indiv_DFC_nolegend.svg"),        # b (nolegend; HGNC-regen)
    "image4.svg":  ("svg", f"{V4}/panel_C_D_DFC_panel_lower_nolegend.svg"),# c (restyled to fig_s11/v2; panel_c_fisher_s11style.R)
    "image6.svg":  ("svg", f"{V4}/panel_C_D_DFC_legend.svg"),             # c legend (restyled)
    "image2.svg":  ("svg", f"{V4}/panel_E_F_DFC_lower.svg"),               # d (HGNC-regen)
    "image12.svg": ("pdf", f"{V4}/panel_e_gsea_table.pdf"),               # e (red highlight + ref superscripts removed; see make_panel_e.py)
    "image14.svg": ("svg", f"{V4}/992_DFC_Schizophrenia_control_bulk_megaset.svg"),  # f (module 992, no red highlight; fig8_specific_modules.R)
}
# Generic render order: c, d, e, f, then the small c-legend on top of c. Panels a/b
# are handled specially below (nolegend dot plots + one shared legend, fig_7/v8 style).
RENDER_ORDER = ["image4.svg", "image2.svg", "image12.svg", "image14.svg", "image6.svg"]

# Single shared legend for panels a/b (replaces the two per-plot legends), placed
# fig_7/v8-style centred below the dot plots. The source SVG has a large transparent
# margin, so crop to its measured content bbox (in the legend's viewBox units).
AB_LEGEND_SRC    = f"{V4}/panel_A_B_legend.svg"
AB_LEGEND_VB     = "31.3 33.7 155.4 23.9"   # content bbox (measured); aspect ~6.5
AB_LEGEND_W_PX   = 188.0                     # rendered width (sized so text ~matches the dot plots)
AB_LEG_BAND_PX   = 32.0                      # vertical band reserved below the plots
AB_LEG_GAP_PX    = 4.0                       # gap between plots and legend

EMU_PER_PX = 9525.0          # 96 dpi
def px(emu): return float(emu) / EMU_PER_PX

NS = {"a":"http://schemas.openxmlformats.org/drawingml/2006/main",
      "p":"http://schemas.openxmlformats.org/presentationml/2006/main",
      "r":"http://schemas.openxmlformats.org/officeDocument/2006/relationships",
      "asvg":"http://schemas.microsoft.com/office/drawing/2016/SVG/main"}

# ---- layout, frozen from v2.pptx (EMU; 914400 per inch) ------------------
# Slide canvas, the panel boxes, and the text shapes drawn over them. Text sizes are
# px (the pptx sz/100 converted at 96/72). Everything else is derived from these.
W_EMU, H_EMU = 7772400, 10058400
W, H = px(W_EMU), px(H_EMU)

pics = {
    "image2.svg":  (4617360, 2711160, 2792880, 1775520),
    "image4.svg":  (1038600, 2644920, 2506680, 2349720),
    "image6.svg":  (3249000, 3217320,  912600,  684000),
    "image8.svg":  (3922920,  603360, 3394080, 1697040),
    "image10.svg": (  90720,  640080, 3437640, 1645560),
    "image12.svg": ( 838800, 5230440, 3657600, 2532960),
    "image14.svg": (5157360, 5111640, 1938600, 3877200),
}

# (paras, x, y, cx, cy, size_px, bold, centred); paras = lines, each a list of
# (text, colour-hex-or-None) runs -- the phrase "shared dCoPA genes" is gold.
# Two template shapes are deliberately absent: its "Fig. 8 v1" placeholder title
# (TITLE is drawn above instead) and a static "r = 0.82" (panel a now draws that as
# a real bracket, from fig8_panelA_bracket.R).
texts = [
    ([[("# of significant dCoPA modules", "000000")], [(" with lower expression in SCZ", "000000")]],
     1918080, 377640, 1229040, 272520, 8.0, False, True),
    ([[("a", None)]], 1152000, 320400, 250200, 242280, 13.3333, True, False),
    ([[("b", None)]], 4896000, 320400, 257400, 242280, 13.3333, True, False),
    ([[("# of significant dCoPA modules", "000000")], [(" with higher expression in SCZ", "000000")]],
     5625720, 378000, 1229040, 272520, 8.0, False, True),
    ([[("d", None)]], 4860360, 2300400, 257760, 242280, 13.3333, True, False),
    ([[("c", None)]], 1152360, 2300400, 250560, 242280, 13.3333, True, False),
    ([[("Cell type overlap of ", "000000"), ("shared dCoPA genes", "b8860b"), (" ", "000000")],
      [("with lower expression in SCZ", "000000")]],
     5416200, 2394000, 1589040, 272520, 8.0, False, True),
    ([[("Significance of overlap of ", "000000"), ("shared dCoPA genes", "b8860b"), (" ", "000000")],
      [("with lower expression in SCZ", "000000")]],
     1609200, 2394360, 1788120, 272520, 8.0, False, True),
    ([[("e", None)]], 1152720, 5000400, 250200, 242280, 13.3333, True, False),
    ([[("Summary of ", "000000"), ("shared dCoPA genes", "b8860b")]],
     1984320, 5058360, 1326600, 181080, 8.0, False, True),
    ([[("f", None)]], 4895280, 5038200, 222120, 242280, 13.3333, True, False),
]

# ---- panel e: PDF -> content-cropped vector SVG --------------------------
subprocess.run(["inkscape", MEDIA_SRC["image12.svg"][1],
                "--export-type=svg", "--export-area-drawing",
                f"--export-filename={PANEL_E_SVG}"],
               check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
SRC_PATH = {m: (PANEL_E_SVG if m == "image12.svg" else MEDIA_SRC[m][1])
            for m in MEDIA_SRC}

# ---- SVG inlining (keeps panels vector; prefix ids to avoid cross-panel clashes)
def _idprefix(s, p):
    s = re.sub(r'\bid="([^"]+)"',          lambda m: f'id="{p}{m.group(1)}"', s)
    s = re.sub(r"\bid='([^']+)'",          lambda m: f"id='{p}{m.group(1)}'", s)
    s = re.sub(r'xlink:href="#([^"]+)"',   lambda m: f'xlink:href="#{p}{m.group(1)}"', s)
    s = re.sub(r"xlink:href='#([^']+)'",   lambda m: f"xlink:href='#{p}{m.group(1)}'", s)
    s = re.sub(r'(?<!xlink:)\bhref="#([^"]+)"', lambda m: f'href="#{p}{m.group(1)}"', s)
    s = re.sub(r'url\(#([^)]+)\)',         lambda m: f'url(#{p}{m.group(1)})', s)
    return s

def inline(path, prefix, x, y, w, h, par="xMidYMid meet", vb_override=None):
    s = open(path, encoding="utf-8").read()
    s = re.sub(r'<\?xml[^>]*\?>', '', s)
    s = re.sub(r'<!DOCTYPE[^>]*>', '', s).strip()
    head = re.search(r'<svg\b([^>]*)>', s).group(1)
    # vb_override (a "minx miny w h" string) crops to a sub-rectangle of the source.
    vb = vb_override or re.search(r'viewBox=["\']([^"\']+)["\']', head).group(1)
    # keep the source's xmlns:* declarations (Inkscape SVGs use rdf/sodipodi/
    # inkscape prefixes inside; dropping their decls breaks the XML parse).
    xmlns = " ".join(re.findall(r'xmlns:[\w-]+=["\'][^"\']*["\']', head))
    overflow = "hidden" if vb_override else "visible"   # clip only when cropping
    s = _idprefix(s, prefix)
    s = re.sub(r'<svg\b[^>]*>',
               f'<svg {xmlns} x="{x:.2f}" y="{y:.2f}" width="{w:.2f}" height="{h:.2f}" '
               f'viewBox="{vb}" preserveAspectRatio="{par}" overflow="{overflow}">',
               s, count=1)
    return s

def esc(t): return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
FAMILY = "Arial, 'Liberation Sans', sans-serif"

# Figure title, drawn bold and spanning the full width across the top (replaces
# the "Fig. 8 v1" placeholder that lives in the template).
TITLE = ("Fig. 8 | dCoPA reveals reproducible downregulation of modular gene "
         "activity in SST inhibitory neurons, VIP inhibitory neurons, and L6b "
         "excitatory neurons in schizophrenia (SCZ)")
TITLE_SIZE   = 13.0
TITLE_MARGIN = 8.0    # px inset from each side

# Layout nudges (px). CONTENT_DY shifts every panel/letter/subtitle down (away from
# the title) without moving the title. EF_EXTRA_DY drops panel e's letter + subtitle
# (and panel f's letter, snapped to e's row) further down toward the table.
CONTENT_DY   = 15.0
EF_EXTRA_DY  = 12.0
LETTER_DY    = 7.0    # drop the a–d letters so they line up with their 2-line subtitles

def wrap(text, max_w, size, bold):
    """Greedy word-wrap to max_w px using an approximate Arial glyph width."""
    cw = size * (0.58 if bold else 0.52)
    lines, cur = [], ""
    for word in text.split():
        trial = word if not cur else f"{cur} {word}"
        if not cur or len(trial) * cw <= max_w:
            cur = trial
        else:
            lines.append(cur); cur = word
    if cur:
        lines.append(cur)
    return lines

parts = [f'<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" '
         f'width="{W:.2f}" height="{H:.2f}" viewBox="0 0 {W:.2f} {H:.2f}">',
         f'<rect width="{W:.2f}" height="{H:.2f}" fill="white"/>']

# figure title: bold, flush-left, wrapped across the full top width. Drawn outside
# the content group so it stays put while everything else shifts down (CONTENT_DY).
for i, line in enumerate(wrap(TITLE, W - 2 * TITLE_MARGIN, TITLE_SIZE, bold=True)):
    ty = 4 + TITLE_SIZE * 0.95 + i * TITLE_SIZE * 1.25
    parts.append(f'<text x="{TITLE_MARGIN:.1f}" y="{ty:.1f}" text-anchor="start" '
                 f'font-family="{FAMILY}" font-size="{TITLE_SIZE:.1f}" font-weight="bold">{esc(line)}</text>')

# everything below the title is shifted down by CONTENT_DY (away from the title)
parts.append(f'<g transform="translate(0,{CONTENT_DY:.1f})">')

# panels (vector), in render order
for media in RENDER_ORDER:
    x, y, w, h = pics[media]
    parts.append(inline(SRC_PATH[media], media.replace(".", "_") + "_",
                         px(x), px(y), px(w), px(h)))

# panels a/b: nolegend dot plots, top-anchored, with one shared legend centred
# below them (fig_7/v8 style). Reserve a band at the bottom of the a/b row for the
# legend; the two dot matrices occupy the area above it.
ax, ay, aw, ah = pics["image10.svg"]
bx, by, bw, bh = pics["image8.svg"]
ab_top_px    = px(min(ay, by))
ab_bottom_px = px(max(ay + ah, by + bh))
matrix_h_px  = (ab_bottom_px - ab_top_px) - AB_LEG_BAND_PX - AB_LEG_GAP_PX
for media in ("image10.svg", "image8.svg"):
    mx, my, mw, mh = pics[media]
    parts.append(inline(SRC_PATH[media], media.replace(".", "_") + "_",
                         px(mx), ab_top_px, px(mw), matrix_h_px, par="xMidYMin meet"))

leg_h_px = AB_LEGEND_W_PX * 23.9 / 155.4                          # keep content aspect
leg_cx   = (px(ax + aw / 2) + px(bx + bw / 2)) / 2                # midway between a & b
leg_x    = leg_cx - AB_LEGEND_W_PX / 2
leg_y    = ab_top_px + matrix_h_px + AB_LEG_GAP_PX
parts.append(inline(AB_LEGEND_SRC, "ablegend_", leg_x, leg_y,
                    AB_LEGEND_W_PX, leg_h_px, vb_override=AB_LEGEND_VB))

# letters (a–f) and per-panel titles, drawn as Arial, preserving the template's
# per-run colours (the "shared dCoPA genes" phrase is gold in the c/d/e subtitles).
# Explicit paragraph lines are stacked; centred shapes are centred on their box.
def plain_lines(paras):
    return ["".join(t for t, _ in ln) for ln in paras]
e_letter_y = next((t[2] for t in texts if plain_lines(t[0]) == ["e"]), None)  # EMU y of "e"
for paras, x, y, cx, cy, sz, bold, centred in texts:
    plain = plain_lines(paras)
    # panel e letter/subtitle + panel f letter: drop further toward the table, with
    # f snapped to e's row so the two letters align.
    is_e_letter = plain == ["e"]
    is_f_letter = plain == ["f"]
    is_e_title  = plain and plain[0].startswith("Summary of shared")
    if is_e_letter or is_f_letter or is_e_title:
        extra = EF_EXTRA_DY
    elif plain in (["a"], ["b"], ["c"], ["d"]):
        extra = LETTER_DY          # line a–d letters up with their subtitles
    else:
        extra = 0.0
    y_src = e_letter_y if (is_f_letter and e_letter_y is not None) else y  # align f with e
    weight = ' font-weight="bold"' if bold else ''
    if centred:
        tx, anchor = px(x + cx / 2), "middle"
    else:
        tx, anchor = px(x) + 2, "start"
    for i, line in enumerate(paras):
        ty = px(y_src) + extra + sz * 0.95 + i * sz * 1.2   # box-top -> baseline, then line-height
        spans = ""
        for rt, rc in line:
            if not rt:
                continue
            fill = f' fill="#{rc}"' if (rc and rc.lower() != "000000") else ""
            spans += f'<tspan{fill}>{esc(rt)}</tspan>'
        parts.append(f'<text x="{tx:.1f}" y="{ty:.1f}" text-anchor="{anchor}" '
                     f'xml:space="preserve" font-family="{FAMILY}" '
                     f'font-size="{sz:.1f}"{weight}>{spans}</text>')

parts.append("</g>")
parts.append("</svg>")
open(MASTER_SVG, "w", encoding="utf-8").write("\n".join(parts))   # intermediate only

# vector PDF (deliverable)
subprocess.run(["rsvg-convert", "-f", "pdf", MASTER_SVG, "-o", OUT_PDF], check=True)
print("wrote:", OUT_PDF, f"({W:.0f}x{H:.0f}px)")

# raster PNG (deliverable). The master SVG is sized in px at 96 dpi, so rsvg's
# --dpi is a no-op; scale the pixel dimensions to render at 300 dpi instead.
PNG_DPI = 300
png_w, png_h = round(W * PNG_DPI / 96), round(H * PNG_DPI / 96)
subprocess.run(["rsvg-convert", "-f", "png", "-w", str(png_w), "-h", str(png_h),
                MASTER_SVG, "-o", OUT_PNG], check=True)
print("wrote:", OUT_PNG, f"({png_w}x{png_h}px @ {PNG_DPI}dpi)")

# ---- PPTX (deliverable): one slide = the composed figure, embedded as a vector
# SVG with a PNG fallback (the native PowerPoint pattern: a:blip -> PNG,
# asvg:svgBlip -> SVG), so the slide stays vector in PowerPoint while remaining
# openable everywhere. The SVG lives only inside the pptx; no .svg file is emitted.
from pptx import Presentation
from pptx.util import Emu
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(int(W_EMU)), Emu(int(H_EMU))
slide = prs.slides.add_slide(prs.slide_layouts[6])          # blank layout
pic = slide.shapes.add_picture(OUT_PNG, 0, 0, Emu(int(W_EMU)), Emu(int(H_EMU)))

# add the master SVG as a media part related from the slide, then point an
# asvg:svgBlip at it inside the picture's a:blip (alongside the PNG fallback).
spart = slide.part
svg_partname = spart.package.next_partname("/ppt/media/image%d.svg")
svg_part = Part(svg_partname, "image/svg+xml",
                spart.package, open(MASTER_SVG, "rb").read())
IMAGE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
svg_rid = spart.relate_to(svg_part, IMAGE_REL)

blip = pic._element.blipFill.find(qn("a:blip"))
blip.append(parse_xml(
    '<a:extLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">'
    '<asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    f'r:embed="{svg_rid}"/></a:ext></a:extLst>'))

prs.save(OUT_PPTX)
for tmp in (MASTER_SVG, PANEL_E_SVG):   # intermediates; not deliverables
    os.remove(tmp)
print("wrote:", OUT_PPTX)
