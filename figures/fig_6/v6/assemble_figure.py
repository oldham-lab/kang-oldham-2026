#!/usr/bin/env python3
"""
Assemble Figure 6 (v6) from its panels into PDF, PNG and PPTX -- de novo.

Replaces the media-swap assembler, which copied the hand-laid-out fig_6/v5/v5.pptx
and substituted its embedded images. That made a superseded version directory a
build input, and the template was never tracked, so the figure could not be
rebuilt from a clone. fig_6/v5 held nothing else, so it retires with this change.

Same house pattern as fig_1/v3 and fig_3/v4: compute placements once, emit the
PDF as vector, the PNG as a 300-dpi raster of it, and the PPTX with panels
rasterised at 300 dpi and the text live. The placements below ARE the layout --
they were read out of the v6 deliverable the old script produced (v5's boxes plus
the title respan and the 0.16 in panel-A nudge it applied at build time), so no
template is consulted.

The old path exported the PDF by way of LibreOffice from a raster-only copy of
the pptx, so every panel arrived rasterised; here the panels stay vector.

Layout (US Letter portrait, 612 x 792 pt):
  a        decopa_schematic_final.svg              dCoPA schematic, full-width banner
  b/d/f/h  83_MTC_allAD_Con_bulk_megaset.svg       left column
  c/e/g/i  147_MTC_allAD_Con_bulk_megaset.svg      right column

Each module column is ONE image carrying four stacked sub-panels; its four letters
are drawn here, at the sub-panel boundaries.

NB the filenames and the panel titles disagree: fig6_specific_modules.R is run
with target_mods = c(83, 147) and names its outputs after those ids, but the
panels it draws are titled "Module 77" and "Module 139". Pre-existing, and left
alone here -- this script only places whatever those files contain.
"""
import os, subprocess, tempfile

import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt

V6 = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.join(V6, "Kang_Figure_6_v6")
OUT_PDF, OUT_PNG, OUT_PPTX = BASE + ".pdf", BASE + ".png", BASE + ".pptx"

# --- page geometry (points; US Letter portrait) -------------------------------
PW, PH = 612.0, 792.0

TITLE = ("Fig. 6 | Differential CoPA (dCoPA) identifies gene coexpression modules "
         "that are significantly and uniformly dysregulated in specific cell types "
         "in disease")
TITLE_BOX = (7.2, 5.33, 597.6, 50.4)   # x, y, w, h -- spans the page, wraps to 2 lines
TITLE_FS, LET_FS = 12.0, 14.0

# PowerPoint insets its shape text; the coordinates here are shape origins, so the
# same offsets fig_1/v3 and fig_3/v4 needed apply. Measured against the old
# LibreOffice-rendered deliverable: title +7.09/+5.88, letters +7.05/+6.28.
INSET_X, INSET_Y_BOX, INSET_Y_LETTER = 7.2, 5.0, 6.25

SRC = {
    "A":    "decopa_schematic_final.svg",
    "L83":  "83_MTC_allAD_Con_bulk_megaset.svg",
    "R147": "147_MTC_allAD_Con_bulk_megaset.svg",
}

# Placement cells (x, y, w, h). Panel A already carries the 0.16 in downward nudge
# the old script applied so the schematic clears the two-line title.
CELLS = {
    "A":    ( 42.49,  44.76, 524.07, 163.81),
    "L83":  ( 61.20, 221.75, 224.64, 524.15),
    "R147": (344.90, 221.75, 224.64, 524.15),
}

# Panel letters. a labels the schematic; the rest label sub-panels inside the two
# module columns, so their y values are sub-panel boundaries, not cell tops.
LETTERS = {
    "a": ( 33.85,  45.22),
    "b": ( 34.02, 215.15), "c": (317.48, 215.15),
    "d": ( 34.02, 331.37), "e": (317.51, 331.37),
    "f": ( 34.02, 393.73), "g": (317.48, 393.73),
    "h": ( 34.02, 589.32), "i": (317.48, 589.32),
}


def stage(tmp):
    """Panel SVG -> PDF. Inkscape rather than rsvg: rsvg renders the schematic's
    coloured boxplots and legend swatches solid black (same trap as fig_3/v4)."""
    out = {}
    for k, name in SRC.items():
        p = os.path.join(tmp, f"{k}.pdf")
        subprocess.run(["inkscape", os.path.join(V6, name), "--export-type=pdf",
                        f"--export-filename={p}"], check=True, capture_output=True)
        out[k] = p
    return out


def content_bbox(src_pdf):
    page = fitz.open(src_pdf)[0]
    r = fitz.Rect(1e9, 1e9, -1e9, -1e9)
    for b in page.get_text("blocks"):
        r |= fitz.Rect(b[:4])
    for d in page.get_drawings():
        r |= d["rect"]

    # get_drawings() does not descend into Form XObjects, and Inkscape puts content
    # inside them -- fig_4/v6's panel S had its whole column dendrogram clipped away
    # by trimming to the under-reported box. Union with a raster ink scan, which sees
    # whatever the page actually paints; the union can only grow the box, never crop.
    # Applied as a SAFETY NET, not a co-equal source: the scan rounds outward by up
    # to a pixel and antialiasing bleeds past the true vector edge, so taking it
    # verbatim would nudge every box by a fraction of a point and shift figures that
    # are already published. Only genuinely missed content -- beyond SCAN_TOL -- wins.
    SCAN_DPI, SCAN_TOL = 144, 1.0
    pm = page.get_pixmap(dpi=SCAN_DPI)
    a = np.frombuffer(pm.samples, dtype=np.uint8).reshape(pm.height, pm.width, pm.n)[:, :, :3]
    ys, xs = np.where((a < 250).any(2))
    if len(ys):
        k = 72.0 / SCAN_DPI
        ink = fitz.Rect(page.rect.x0 + xs.min() * k, page.rect.y0 + ys.min() * k,
                        page.rect.x0 + (xs.max() + 1) * k, page.rect.y0 + (ys.max() + 1) * k)
        if ink.x0 < r.x0 - SCAN_TOL: r.x0 = ink.x0
        if ink.y0 < r.y0 - SCAN_TOL: r.y0 = ink.y0
        if ink.x1 > r.x1 + SCAN_TOL: r.x1 = ink.x1
        if ink.y1 > r.y1 + SCAN_TOL: r.y1 = ink.y1

    if r.is_empty or r.x1 < r.x0:
        return page.rect
    return r & page.rect


def fit(bb, cx, cy, cw, ch):
    """Aspect-preserving fit, centred in the cell -- matching the old script, which
    centred each swapped picture in its box."""
    s = min(cw / bb.width, ch / bb.height)
    dw, dh = bb.width * s, bb.height * s
    return fitz.Rect(cx + (cw - dw) / 2, cy + (ch - dh) / 2,
                     cx + (cw - dw) / 2 + dw, cy + (ch - dh) / 2 + dh)


def compute_layout(pdf):
    panels = []
    for k, cell in CELLS.items():
        bb = content_bbox(pdf[k])
        panels.append((k, pdf[k], bb, fit(bb, *cell)))
    return {"panels": panels}


def render_pdf(lay):
    doc = fitz.open()
    page = doc.new_page(width=PW, height=PH)
    x, y, w, h = TITLE_BOX
    page.insert_textbox(fitz.Rect(x + INSET_X, y + INSET_Y_BOX, x + w, y + INSET_Y_BOX + h),
                        TITLE, fontname="hebo", fontsize=TITLE_FS,
                        align=fitz.TEXT_ALIGN_LEFT, lineheight=1.15)
    for _, src, bb, rect in lay["panels"]:
        page.show_pdf_page(rect, fitz.open(src), 0, clip=bb)
    for ch, (lx, ly) in LETTERS.items():
        page.insert_text((lx + INSET_X, ly + INSET_Y_LETTER + LET_FS), ch,
                         fontname="hebo", fontsize=LET_FS)
    doc.save(OUT_PDF, deflate=True)
    return doc


def render_png(doc):
    doc[0].get_pixmap(dpi=300).save(OUT_PNG)


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, bold=False):
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    return r


def render_pptx(lay, tmp):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(PW), Pt(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for key, src, bb, rect in lay["panels"]:
        png = os.path.join(tmp, f"{key}_raster.png")
        fitz.open(src)[0].get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72),
                                     clip=bb).save(png)
        slide.shapes.add_picture(png, Pt(rect.x0), Pt(rect.y0),
                                 Pt(rect.width), Pt(rect.height))

    x, y, w, h = TITLE_BOX
    _run(_box(slide, x + INSET_X, y + INSET_Y_BOX, w, h).paragraphs[0],
         TITLE, TITLE_FS, bold=True)
    for ch, (lx, ly) in LETTERS.items():
        _run(_box(slide, lx + INSET_X, ly + INSET_Y_BOX, 24.0,
                  LET_FS * 1.7).paragraphs[0], ch, LET_FS, bold=True)

    prs.save(OUT_PPTX)


def main():
    with tempfile.TemporaryDirectory() as tmp:
        pdf = stage(tmp)
        lay = compute_layout(pdf)
        doc = render_pdf(lay)
        render_png(doc)
        render_pptx(lay, tmp)
    for p in (OUT_PDF, OUT_PNG, OUT_PPTX):
        print(f"wrote {p}  ({os.path.getsize(p)/1e6:.2f} MB)")


if __name__ == "__main__":
    main()
