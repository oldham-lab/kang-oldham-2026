#!/usr/bin/env python3
"""
Assemble Figure 3 (v4) from its panels into PDF, PNG and PPTX -- de novo.

Replaces the media-swap approach of assemble_figure.py, which copied the
hand-laid-out fig_3/v3/v3.pptx and substituted its embedded images. That made a
superseded version directory a build input, and the template itself was never
tracked, so the figure could not be rebuilt from a clone.

The layout is now expressed here, in points, following the house pattern used by
fig_7/v8 and the fig_s3/s4/s5 assemblers: compute placements once, then emit the
PDF as vector, the PNG as a 300-dpi raster of that PDF, and the PPTX with the
panels rasterised at 300 dpi and the text live/editable.

Geometry was lifted verbatim from v3.pptx (converted EMU -> pt) so this produces
the same approved layout; the numbers are now data in this file rather than a
binary input. The 12 "r = ..." labels of the v3 deck are intentionally absent --
panels E and F carry their brackets and r labels baked into their SVGs.

Layout (US Letter portrait, 612 x 792 pt), by panel letter:
  a  layout_final.svg                       (top schematic banner)
  b  panel_B_bc.svg
  c  panel_C.svg
  d  panel_D.svg
  e  panel_E_native_log_with_brackets.svg
  f  panel_F_REI_with_brackets.svg
panel_B_seed.svg exists in v4/ but is not placed, matching the v3 layout.
"""
import os, subprocess, tempfile

import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt

V4 = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.join(V4, "Kang_Figure_3_v4")
OUT_PDF, OUT_PNG, OUT_PPTX = BASE + ".pdf", BASE + ".png", BASE + ".pptx"

# --- page geometry (points; US Letter portrait) -------------------------------
PW, PH = 612.0, 792.0
MARG = 12.7

TITLE = ("Fig. 3 | Covariation Projection Analysis (CoPA) reveals the "
         "cellular origins of bulk gene coexpression modules")
TITLE_TOP, TITLE_H, TITLE_FS = 11.0, 40.0, 12.0   # 30 overflows: insert_textbox returns <0 and draws nothing
LET_FS, CAP_FS = 14.0, 7.0

# The CELLS/LETTERS/CAPTIONS coordinates below are PowerPoint *shape* origins
# taken from v3.pptx. PowerPoint insets text inside its shapes (0.1 in left,
# 0.05 in top by default), so text drawn at the raw origin sits up and to the
# left of where the approved figure has it. These are the measured offsets
# between this script's output and Kang_Figure_3_v4.pdf, and they were uniform
# across all six letters and all three captions.
INSET_X, INSET_Y_BOX, INSET_Y_LETTER = 7.2, 5.0, 7.3

# panel source SVGs
SRC = {
    "a": "layout_final.svg",
    "b": "panel_B_bc.svg",
    "c": "panel_C.svg",
    "d": "panel_D.svg",
    "e": "panel_E_native_log_with_brackets.svg",
    "f": "panel_F_REI_with_brackets.svg",
}

# Placement cells (x, y, w, h) -- the v3 picture frames, EMU/12700.
# Each panel is aspect-fitted inside its cell, so a regenerated panel whose
# proportions drifted is letterboxed rather than stretched.
CELLS = {
    "a": (95.8,  55.2, 424.2, 105.6),
    "b": ( 8.1, 170.6, 182.9, 133.2),
    "c": (198.0, 194.3, 197.9, 107.9),
    "d": (387.6, 194.3, 224.3,  57.3),
    "e": (  1.0, 344.6, 284.2, 236.8),
    "f": (304.0, 343.6, 299.5, 239.0),
}

# panel letters: (x, y_baseline_box_top)
LETTERS = {
    "a": ( 63.7,  39.3),
    "b": (  9.5, 164.4),
    "c": (199.4, 164.4),
    "d": (389.3, 164.4),
    "e": (  9.5, 325.9),
    "f": (324.1, 328.8),
}

# free-standing captions: (text, x, y, w)
CAPTIONS = [
    ("Pairwise correlations of module genes \nin constituent datasets",
                                              239.8, 171.4, 133.7),
    ("Gene set enrichment analysis",          452.8, 173.1, 105.8),
    ("Gene coexpression module projections on snRNA-seq datasets",
                                              190.8, 321.7, 210.0),
]


def stage(tmp):
    """SVG -> PDF so panels can be placed as vector.

    Inkscape, not rsvg-convert: rsvg renders the panel A schematic's colour
    correlation matrices as solid black. The media-swap script hit the same
    thing and rasterised that panel with Inkscape for the same reason.
    """
    out = {}
    for k, name in SRC.items():
        p = os.path.join(tmp, f"{k}.pdf")
        subprocess.run(["inkscape", os.path.join(V4, name),
                        "--export-type=pdf", f"--export-filename={p}"],
                       check=True, capture_output=True)
        out[k] = p
    return out


def content_bbox(src_pdf):
    """Ink bounds of a panel, so surrounding whitespace does not eat the cell."""
    page = fitz.open(src_pdf)[0]
    r = fitz.Rect(1e9, 1e9, -1e9, -1e9)
    for b in page.get_text("blocks"):
        r |= fitz.Rect(b[:4])
    for d in page.get_drawings():
        r |= d["rect"]

    # get_drawings() does not descend into Form XObjects, and Inkscape puts content
    # inside them -- fig_4/v6's panel S had its whole column dendrogram clipped away
    # by trimming to the under-reported box. Union with a raster ink scan, which sees
    # whatever the page actually paints; the union can only grow the box, never crop.
    # Applied as a SAFETY NET, not a co-equal source: the scan rounds outward by up
    # to a pixel and antialiasing bleeds past the true vector edge, so taking it
    # verbatim would nudge every box by a fraction of a point and shift figures that
    # are already published. Only genuinely missed content -- beyond SCAN_TOL -- wins.
    SCAN_DPI, SCAN_TOL = 144, 1.0
    pm = page.get_pixmap(dpi=SCAN_DPI)
    a = np.frombuffer(pm.samples, dtype=np.uint8).reshape(pm.height, pm.width, pm.n)[:, :, :3]
    ys, xs = np.where((a < 250).any(2))
    if len(ys):
        k = 72.0 / SCAN_DPI
        ink = fitz.Rect(page.rect.x0 + xs.min() * k, page.rect.y0 + ys.min() * k,
                        page.rect.x0 + (xs.max() + 1) * k, page.rect.y0 + (ys.max() + 1) * k)
        if ink.x0 < r.x0 - SCAN_TOL: r.x0 = ink.x0
        if ink.y0 < r.y0 - SCAN_TOL: r.y0 = ink.y0
        if ink.x1 > r.x1 + SCAN_TOL: r.x1 = ink.x1
        if ink.y1 > r.y1 + SCAN_TOL: r.y1 = ink.y1

    if r.is_empty or r.x1 < r.x0:
        return page.rect
    return r & page.rect


def fit(bb, cx, cy, cw, ch):
    """Aspect-preserving fit of bb into the cell, centred."""
    s = min(cw / bb.width, ch / bb.height)
    dw, dh = bb.width * s, bb.height * s
    return fitz.Rect(cx + (cw - dw) / 2, cy + (ch - dh) / 2,
                     cx + (cw - dw) / 2 + dw, cy + (ch - dh) / 2 + dh)


def compute_layout(pdf):
    panels = []
    for k, cell in CELLS.items():
        bb = content_bbox(pdf[k])
        panels.append((k, pdf[k], bb, fit(bb, *cell)))
    return {"panels": panels}


def render_pdf(lay):
    doc = fitz.open()
    page = doc.new_page(width=PW, height=PH)
    page.insert_textbox(fitz.Rect(MARG + INSET_X, TITLE_TOP + INSET_Y_BOX,
                                  PW - MARG, TITLE_TOP + INSET_Y_BOX + TITLE_H),
                        TITLE, fontname="hebo", fontsize=TITLE_FS,
                        align=fitz.TEXT_ALIGN_LEFT, lineheight=1.15)
    for _, src, bb, rect in lay["panels"]:
        page.show_pdf_page(rect, fitz.open(src), 0, clip=bb)
    for ch, (x, y) in LETTERS.items():
        page.insert_text((x + INSET_X, y + INSET_Y_LETTER + LET_FS), ch,
                         fontname="hebo", fontsize=LET_FS)
    for text, x, y, w in CAPTIONS:
        page.insert_textbox(fitz.Rect(x + INSET_X, y + INSET_Y_BOX,
                                      x + INSET_X + w, y + INSET_Y_BOX + 3 * CAP_FS), text,
                            fontname="helv", fontsize=CAP_FS, lineheight=1.15)
    doc.save(OUT_PDF, deflate=True)
    return doc


def render_png(doc):
    doc[0].get_pixmap(dpi=300).save(OUT_PNG)


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, bold=False):
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    return r


def render_pptx(lay, tmp):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(PW), Pt(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    for key, src, bb, rect in lay["panels"]:
        png = os.path.join(tmp, f"{key}_raster.png")
        fitz.open(src)[0].get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72),
                                     clip=bb).save(png)
        slide.shapes.add_picture(png, Pt(rect.x0), Pt(rect.y0),
                                 Pt(rect.width), Pt(rect.height))

    _run(_box(slide, MARG + INSET_X, TITLE_TOP + INSET_Y_BOX,
              PW - 2 * MARG, TITLE_H).paragraphs[0], TITLE, TITLE_FS, bold=True)
    for ch, (x, y) in LETTERS.items():
        _run(_box(slide, x + INSET_X, y + INSET_Y_BOX, 24.0,
                  LET_FS * 1.7).paragraphs[0], ch, LET_FS, bold=True)
    for text, x, y, w in CAPTIONS:
        _run(_box(slide, x + INSET_X, y + INSET_Y_BOX, w,
                  3 * CAP_FS).paragraphs[0], text, CAP_FS)

    prs.save(OUT_PPTX)


def main():
    with tempfile.TemporaryDirectory() as tmp:
        pdf = stage(tmp)
        lay = compute_layout(pdf)
        doc = render_pdf(lay)
        render_png(doc)
        render_pptx(lay, tmp)
    for p in (OUT_PDF, OUT_PNG, OUT_PPTX):
        print(f"wrote {p}  ({os.path.getsize(p)/1e6:.2f} MB)")


if __name__ == "__main__":
    main()
