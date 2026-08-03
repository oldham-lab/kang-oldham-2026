#!/usr/bin/env python3
"""
Assemble Figure S1 from its individual panels into PDF, PNG and PPTX.

fig_s1 compares the proportion of nuclei assigned to each cell class, subclass and
supertype in Jorstad et al. 2023 versus Gabitto et al. 2024. Layout (2 rows x 1 col),
each row a 3-facet strip produced by fig_s1_v3.R:
  a-c : DFC (panel_AtoC)
  d-f : MTG (panel_DtoF)

The geometry is taken from the hand-built fig_s1_v3.pptx in this folder, with two
deliberate changes: a title is added in the Fig. 1 / Fig. S2 house style (Arial 12 pt
bold at the top-left corner), and the panel stack is shifted down by TITLE_DY to clear
it. Letter tops in the hand-built deck drift by ~2 pt within a row (18.0/20.1/20.1);
they are normalised to the row's first letter here.

Mirrors the other assemble_figure.py scripts: the layout is computed once; the PDF
places panels as vector, the PNG is a 300-dpi raster of the PDF, and the PPTX places
300-dpi panel images at the same positions with live (editable) title and letters.
The panel letters are NOT baked into the ggplot output -- they are native text.
"""
import os, tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

FIG  = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_s1/v3")
BASE = os.path.join(FIG, "Kang_Figure_S1_v3")
OUT_PDF, OUT_PNG, OUT_PPTX = BASE + ".pdf", BASE + ".png", BASE + ".pptx"

TITLE = ("Fig. S1 | Cell class and subclass proportions are more highly conserved "
         "than supertype proportions")

# --- page + layout geometry (points; from fig_s1_v3.pptx) --------------------------
PW, PH = 612.0, 792.0

TITLE_X, TITLE_Y, TITLE_W, TITLE_H, TITLE_FS = 7.2, 2.8, 597.6, 20.0, 12.0
TITLE_DY = 11.0   # push the panel stack down to clear the title

LET_FS = 18.0
LET_X = (30.3, 216.0, 405.9)          # shared by both rows
LET_INSET_X, LET_BASELINE_DY = 7.2, 18.6   # textbox origin -> glyph origin, 18 pt Arial

# each row: (source pdf, letters, picture rect, letter-box top)
ROWS = [
    (f"{FIG}/panel_AtoC.pdf", ("a", "b", "c"), (2.9, 41.5, 594.0, 208.1), 18.0),
    (f"{FIG}/panel_DtoF.pdf", ("d", "e", "f"), (2.9, 266.6, 593.9, 207.8), 244.8),
]


def compute_layout():
    """Return placements shared by every output format."""
    panels, letters = [], []
    for src, chars, (x, y, w, h), lt in ROWS:
        panels.append((src, fitz.Rect(x, y + TITLE_DY, x + w, y + TITLE_DY + h)))
        for ch, lx in zip(chars, LET_X):
            letters.append((ch, lx, lt + TITLE_DY))
    return {"panels": panels, "letters": letters}


# ---- PDF (vector) ------------------------------------------------------------
def render_pdf(lay):
    doc = fitz.open()
    page = doc.new_page(width=PW, height=PH)
    page.insert_textbox(fitz.Rect(TITLE_X, TITLE_Y, TITLE_X + TITLE_W, TITLE_Y + TITLE_H),
                        TITLE, fontname="hebo", fontsize=TITLE_FS,
                        align=fitz.TEXT_ALIGN_LEFT, lineheight=1.15)
    for src, rect in lay["panels"]:
        page.show_pdf_page(rect, fitz.open(src), 0)
    for ch, x, y in lay["letters"]:
        page.insert_text((x + LET_INSET_X, y + LET_BASELINE_DY), ch,
                         fontname="hebo", fontsize=LET_FS)
    doc.save(OUT_PDF, deflate=True)
    return doc


def render_png(doc):
    doc[0].get_pixmap(dpi=300).save(OUT_PNG)


# ---- PPTX (panels rasterised at 300 dpi, text editable) ----------------------
def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, bold=False):
    r = p.add_run(); r.text = text
    r.font.name = "Arial"; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = RGBColor(0, 0, 0)
    return r


def render_pptx(lay, tmp):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(PW), Pt(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    for i, (src, rect) in enumerate(lay["panels"]):
        png = os.path.join(tmp, f"row{i}_raster.png")
        fitz.open(src)[0].get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72)).save(png)
        slide.shapes.add_picture(png, Pt(rect.x0), Pt(rect.y0),
                                 Pt(rect.width), Pt(rect.height))

    tf = _box(slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H)
    _run(tf.paragraphs[0], TITLE, TITLE_FS, bold=True)

    for ch, x, y in lay["letters"]:
        tf = _box(slide, x + LET_INSET_X, y, LET_FS * 1.6, LET_FS * 1.6)
        _run(tf.paragraphs[0], ch, LET_FS, bold=True)

    prs.save(OUT_PPTX)


def main():
    tmp = tempfile.mkdtemp(prefix="figs1_assemble_")
    lay = compute_layout()
    doc = render_pdf(lay)
    render_png(doc)
    render_pptx(lay, tmp)
    for p in (OUT_PDF, OUT_PNG, OUT_PPTX):
        print("Saved:", p)


if __name__ == "__main__":
    main()
