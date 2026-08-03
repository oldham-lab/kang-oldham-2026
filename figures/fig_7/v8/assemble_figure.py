#!/usr/bin/env python3
"""
Assemble Figure 7 (v8) from its individual panels into PDF, PNG and PPTX.

Layout (3 rows x 2 cols), MTG on the left, DFC on the right:
  a / b : dCoPA dot plots  (panel_A_indiv_{MTG,DFC}_nolegend_bracket, v7.1) + shared legend
  c / d : cell-type overlap UpSet plots (panel_E_F_{MTG,DFC}, v7.1)
  e / f : shared-gene summary tables (panel_G_gsea_summary_table_{mtg,dfc}_fdr, v7.2)

The layout is computed once; the PDF places panels as vector (SVG panels converted via
rsvg-convert), the PNG is a 300-dpi raster of the PDF, and the PPTX places 300-dpi panel
images at the same positions with live (editable) title / letters / subtitles. Panel
letters, per-panel subtitles and the title are native text; "shared dCoPA genes" is gold.
"""
import os, subprocess, tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIG  = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_7")
BASE = os.path.join(FIG, "v8", "Kang_Figure_7_v8")
OUT_PDF, OUT_PNG, OUT_PPTX = BASE + ".pdf", BASE + ".png", BASE + ".pptx"

SRC = {
    "a":      (f"{FIG}/v8/panel_A_indiv_MTG_nolegend.pdf", "pdf"),
    "b":      (f"{FIG}/v8/panel_A_indiv_DFC_nolegend.pdf", "pdf"),
    "legend": (f"{FIG}/v8/panel_A_B_legend.svg", "svg"),
    "c":      (f"{FIG}/v8/panel_E_F_MTG.svg", "svg"),
    "d":      (f"{FIG}/v8/panel_E_F_DFC.svg", "svg"),
    "e":      (f"{FIG}/v8/panel_G_gsea_summary_table_mtg_fdr.pdf", "pdf"),
    "f":      (f"{FIG}/v8/panel_G_gsea_summary_table_dfc_fdr.pdf", "pdf"),
}

TITLE = ("Fig. 7 | dCoPA reveals reproducible downregulation of modular gene activity "
         "in deep-layer intratelencephalic excitatory neurons, PVALB inhibitory neurons, "
         "and LAMP5 inhibitory neurons in Alzheimer's disease (AD)")

# per-panel subtitles, split (prefix, gold-phrase, suffix); gold phrase "" => plain a/b
SUB = {
    "a": ("# of significant dCoPA modules with lower expression in AD MTG", "", ""),
    "b": ("# of significant dCoPA modules with lower expression in AD DFC", "", ""),
    "c": ("Cell type overlap of ", "shared dCoPA genes", " (MTG)"),
    "d": ("Cell type overlap of ", "shared dCoPA genes", " (DFC)"),
    "e": ("Summary of ", "shared dCoPA genes", " (MTG)"),
    "f": ("Summary of ", "shared dCoPA genes", " (DFC)"),
}
GOLD = (184, 134, 11)   # #B8860B, matches the panel axis labels

# --- page + layout geometry (points; US Letter portrait) ----------------------
PW, PH = 612.0, 792.0
MARG, GUT = 22.0, 12.0
COLW = (PW - 2 * MARG - GUT) / 2.0
XL = MARG
XR = MARG + COLW + GUT

TITLE_TOP, TITLE_H, TITLE_FS = 16.0, 52.0, 10.0
SUB_AB_FS, SUB_1L_FS, LET_FS = 7.5, 8.0, 12.0
HDR_AB, HDR_1L = 24.0, 15.0     # header band heights (a/b two-line vs c-f one-line)


def stage(tmp):
    out = {}
    for k, (path, kind) in SRC.items():
        if kind == "pdf":
            out[k] = path
        else:
            p = os.path.join(tmp, f"{k}.pdf")
            subprocess.run(["rsvg-convert", "-f", "pdf", "-o", p, path], check=True)
            out[k] = p
    return out


def content_bbox(src_pdf):
    page = fitz.open(src_pdf)[0]
    r = fitz.Rect(1e9, 1e9, -1e9, -1e9)
    for b in page.get_text("blocks"):
        r |= fitz.Rect(b[:4])
    for d in page.get_drawings():
        r |= d["rect"]
    if r.is_empty or r.x1 < r.x0:
        return page.rect
    return r & page.rect


def fit(bb, cx, cy, cw, ch, halign="left"):
    s = min(cw / bb.width, ch / bb.height)
    dw, dh = bb.width * s, bb.height * s
    dx = cx + (cw - dw) * {"left": 0, "center": .5, "right": 1}[halign]
    return fitz.Rect(dx, cy, dx + dw, cy + dh)


def compute_layout(pdf):
    """Return dicts of placements shared by every output format."""
    bb = {k: content_bbox(v) for k, v in pdf.items()}
    panels, letters, subs = [], [], []

    def add(key, x, y, cw, ch, halign="left"):
        panels.append((key, pdf[key], bb[key], fit(bb[key], x, y, cw, ch, halign)))

    # row 1: dot plots + legend sized to the dot-plot scale (matched text size)
    y = TITLE_TOP + TITLE_H + 12
    h1 = COLW / max(bb["a"].width / bb["a"].height, bb["b"].width / bb["b"].height)
    for key, x in (("a", XL), ("b", XR)):
        letters.append((key, x, y + 10)); subs.append((key, x, y, COLW))
        add(key, x, y + HDR_AB, COLW, h1)
    y += HDR_AB + h1

    s_panel = COLW / bb["a"].width                 # scale used for the dot plots
    lw = bb["legend"].width * s_panel              # legend at same scale => same text size
    lh = bb["legend"].height * s_panel
    leg_gap = 11   # vertical gap below the dot plots before the shared legend
    add("legend", (PW - lw) / 2, y + leg_gap, lw, lh, "center")
    y += leg_gap + lh

    # row 2: UpSet plots
    y += 18
    h2 = COLW / (bb["c"].width / bb["c"].height)
    for key, x in (("c", XL), ("d", XR)):
        letters.append((key, x, y + 10)); subs.append((key, x, y, COLW))
        add(key, x, y + HDR_1L, COLW, h2)
    y += HDR_1L + h2

    # row 3: summary tables
    y += 18
    avail = PH - MARG - (y + HDR_1L)
    for key, x in (("e", XL), ("f", XR)):
        letters.append((key, x, y + 10)); subs.append((key, x, y, COLW))
        add(key, x, y + HDR_1L, COLW, min(COLW / (bb[key].width / bb[key].height), avail))

    return {"panels": panels, "letters": letters, "subs": subs}


# ---- PDF (vector) ------------------------------------------------------------
def render_pdf(lay):
    doc = fitz.open()
    page = doc.new_page(width=PW, height=PH)
    page.insert_textbox(fitz.Rect(MARG, TITLE_TOP, PW - MARG, TITLE_TOP + TITLE_H),
                        TITLE, fontname="hebo", fontsize=TITLE_FS,
                        align=fitz.TEXT_ALIGN_LEFT, lineheight=1.15)
    for key, src, bb, rect in lay["panels"]:
        page.show_pdf_page(rect, fitz.open(src), 0, clip=bb)
    for ch, x, y in lay["letters"]:
        page.insert_text((x, y), ch, fontname="hebo", fontsize=LET_FS)
    for key, cx, cy, cw in lay["subs"]:
        pre, mid, suf = SUB[key]
        if not mid:
            page.insert_textbox(fitz.Rect(cx, cy, cx + cw, cy + HDR_AB), pre,
                                fontname="helv", fontsize=SUB_AB_FS,
                                align=fitz.TEXT_ALIGN_CENTER, lineheight=1.1)
        else:
            fs = SUB_1L_FS
            wp = fitz.get_text_length(pre, "helv", fs)
            wm = fitz.get_text_length(mid, "helv", fs)
            ws = fitz.get_text_length(suf, "helv", fs)
            x0, yb = cx + (cw - (wp + wm + ws)) / 2, cy + fs + 1
            gold = tuple(c / 255 for c in GOLD)
            page.insert_text((x0, yb), pre, fontname="helv", fontsize=fs)
            page.insert_text((x0 + wp, yb), mid, fontname="helv", fontsize=fs, color=gold)
            page.insert_text((x0 + wp + wm, yb), suf, fontname="helv", fontsize=fs)
    doc.save(OUT_PDF, deflate=True)
    return doc


def render_png(doc):
    doc[0].get_pixmap(dpi=300).save(OUT_PNG)


# ---- PPTX (panels rasterised at 300 dpi, text editable) ----------------------
def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, bold=False, color=None):
    r = p.add_run(); r.text = text
    r.font.name = "Arial"; r.font.size = Pt(size); r.font.bold = bold
    if color:
        r.font.color.rgb = RGBColor(*color)
    return r


def render_pptx(lay, tmp):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(PW), Pt(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    for key, src, bb, rect in lay["panels"]:
        png = os.path.join(tmp, f"{key}_raster.png")
        fitz.open(src)[0].get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72), clip=bb).save(png)
        slide.shapes.add_picture(png, Pt(rect.x0), Pt(rect.y0),
                                 Pt(rect.width), Pt(rect.height))

    tf = _box(slide, MARG, TITLE_TOP, PW - 2 * MARG, TITLE_H)
    _run(tf.paragraphs[0], TITLE, TITLE_FS, bold=True)

    for ch, x, y in lay["letters"]:
        tf = _box(slide, x, y - LET_FS, LET_FS + 4, LET_FS + 4)
        _run(tf.paragraphs[0], ch, LET_FS, bold=True)

    for key, cx, cy, cw in lay["subs"]:
        pre, mid, suf = SUB[key]
        tf = _box(slide, cx, cy, cw, HDR_AB if not mid else HDR_1L)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        if not mid:
            _run(p, pre, SUB_AB_FS)
        else:
            _run(p, pre, SUB_1L_FS)
            _run(p, mid, SUB_1L_FS, color=GOLD)
            _run(p, suf, SUB_1L_FS)

    prs.save(OUT_PPTX)


def main():
    tmp = tempfile.mkdtemp(prefix="fig7_assemble_")
    pdf = stage(tmp)
    lay = compute_layout(pdf)
    doc = render_pdf(lay)
    render_png(doc)
    render_pptx(lay, tmp)
    for p in (OUT_PDF, OUT_PNG, OUT_PPTX):
        print("Saved:", p)


if __name__ == "__main__":
    main()
