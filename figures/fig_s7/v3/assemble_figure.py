#!/usr/bin/env python3
"""
Assemble Figure S7 (v3) from its panel into PDF, PNG and PPTX -- de novo.

Replaces the media-swap assembler, which copied Kang_Figure_S7_v3_template.pptx
and substituted its embedded image. The template was never tracked, so the figure
could not be rebuilt from a clone; the placements below ARE the layout, read out
of the deliverable that script produced.

Same house pattern as fig_1/v3, fig_3/v4, fig_4/v6 and fig_6/v6: compute the
placement once, emit the PDF as vector, the PNG as a 300-dpi raster of it, and the
PPTX with the panel rasterised at 300 dpi and the text live. The old path exported
the PDF through LibreOffice from a raster-only copy of the pptx, so the panel
arrived rasterised; here it stays vector.

S7 is a single panel -- the per-subclass genome-wide expression barplots -- so
there are no panel letters. The subtitle's gene count is NOT a literal: it is
rewritten from common_genes.txt on every build (v3's common gene universe), so the
caption always reports the number of genes actually projected.
"""
import os, subprocess, tempfile

import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt

V3 = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.join(V3, "Kang_Figure_S7_v3")
OUT_PDF, OUT_PNG, OUT_PPTX = BASE + ".pdf", BASE + ".png", BASE + ".pptx"

SRC_PANEL = "panel_1.svg"
COMMON_GENES = os.path.join(V3, "common_genes.txt")

# --- page geometry (points; US Letter portrait) -------------------------------
PW, PH = 612.0, 792.0

TITLE = "Fig. S7 | Genome-wide expression levels vary significantly among neocortical subclasses."
TITLE_BOX = (6.66, 2.49, 528.41, 21.46)
TITLE_FS = 12.0

# Two centred lines; {n} is filled from COMMON_GENES at build time.
SUBTITLE = ("Genome-wide projections (n = {n:,} genes)\n"
            " on normal samples from seven snRNA-seq datasets")
SUBTITLE_BOX = (178.72, 57.69, 258.43, 32.23)
SUBTITLE_FS = 10.5

# This deck's text insets are ~1 pt tighter than the other figures' (measured
# against the old LibreOffice render: title +6.15/+5.92, subtitle +5.61 vertically).
INSET_X, INSET_Y_BOX = 6.15, 5.9
# The subtitle is two lines, so its spacing matters: PyMuPDF's lineheight is not
# exactly fontsize x factor, and 1.2 put the second line ~0.9 pt low against the
# reference render. 1.117 reproduces the 12.61 pt baseline gap.
SUBTITLE_INSET_Y, SUBTITLE_LH = 5.61, 1.117

CELL = (90.03, 89.05, 395.91, 539.88)   # panel placement box


def read_gene_count():
    return sum(1 for line in open(COMMON_GENES) if line.strip())


def stage(tmp):
    """Panel SVG -> PDF. Inkscape rather than rsvg (see fig_3/v4)."""
    p = os.path.join(tmp, "panel.pdf")
    subprocess.run(["inkscape", os.path.join(V3, SRC_PANEL), "--export-type=pdf",
                    f"--export-filename={p}"], check=True, capture_output=True)
    return p


def content_bbox(src_pdf):
    """Ink box of the staged panel."""
    page = fitz.open(src_pdf)[0]
    r = fitz.Rect(1e9, 1e9, -1e9, -1e9)
    for b in page.get_text("blocks"):
        r |= fitz.Rect(b[:4])
    for d in page.get_drawings():
        r |= d["rect"]

    # Applied as a SAFETY NET, not a co-equal source: get_drawings() does not descend
    # into Form XObjects (fig_4/v6's panel S lost its dendrogram that way), but the
    # scan also rounds outward and antialiasing bleeds past the true edge, so only
    # genuinely missed content -- beyond SCAN_TOL -- wins.
    SCAN_DPI, SCAN_TOL = 144, 1.0
    pm = page.get_pixmap(dpi=SCAN_DPI)
    a = np.frombuffer(pm.samples, dtype=np.uint8).reshape(pm.height, pm.width, pm.n)[:, :, :3]
    ys, xs = np.where((a < 250).any(2))
    if len(ys):
        k = 72.0 / SCAN_DPI
        ink = fitz.Rect(page.rect.x0 + xs.min() * k, page.rect.y0 + ys.min() * k,
                        page.rect.x0 + (xs.max() + 1) * k, page.rect.y0 + (ys.max() + 1) * k)
        if ink.x0 < r.x0 - SCAN_TOL: r.x0 = ink.x0
        if ink.y0 < r.y0 - SCAN_TOL: r.y0 = ink.y0
        if ink.x1 > r.x1 + SCAN_TOL: r.x1 = ink.x1
        if ink.y1 > r.y1 + SCAN_TOL: r.y1 = ink.y1

    if r.is_empty or r.x1 < r.x0:
        return page.rect
    return r & page.rect


def fit(bb, cx, cy, cw, ch):
    """Aspect-preserving fit, centred in the cell."""
    s = min(cw / bb.width, ch / bb.height)
    dw, dh = bb.width * s, bb.height * s
    return fitz.Rect(cx + (cw - dw) / 2, cy + (ch - dh) / 2,
                     cx + (cw - dw) / 2 + dw, cy + (ch - dh) / 2 + dh)


def render_pdf(src, bb, rect, subtitle):
    doc = fitz.open()
    page = doc.new_page(width=PW, height=PH)
    x, y, w, h = TITLE_BOX
    page.insert_textbox(fitz.Rect(x + INSET_X, y + INSET_Y_BOX, x + w, y + INSET_Y_BOX + h),
                        TITLE, fontname="hebo", fontsize=TITLE_FS,
                        align=fitz.TEXT_ALIGN_LEFT, lineheight=1.15)
    page.show_pdf_page(rect, fitz.open(src), 0, clip=bb)
    x, y, w, h = SUBTITLE_BOX
    page.insert_textbox(fitz.Rect(x, y + SUBTITLE_INSET_Y, x + w, y + SUBTITLE_INSET_Y + h),
                        subtitle, fontname="helv", fontsize=SUBTITLE_FS,
                        align=fitz.TEXT_ALIGN_CENTER, lineheight=SUBTITLE_LH)
    doc.save(OUT_PDF, deflate=True)
    return doc


def render_png(doc):
    doc[0].get_pixmap(dpi=300).save(OUT_PNG)


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def render_pptx(src, bb, rect, subtitle, tmp):
    from pptx.enum.text import PP_ALIGN
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(PW), Pt(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    png = os.path.join(tmp, "panel_raster.png")
    fitz.open(src)[0].get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72), clip=bb).save(png)
    slide.shapes.add_picture(png, Pt(rect.x0), Pt(rect.y0), Pt(rect.width), Pt(rect.height))

    x, y, w, h = TITLE_BOX
    p = _box(slide, x + INSET_X, y + INSET_Y_BOX, w, h).paragraphs[0]
    r = p.add_run(); r.text = TITLE
    r.font.name, r.font.size, r.font.bold = "Arial", Pt(TITLE_FS), True

    x, y, w, h = SUBTITLE_BOX
    tf = _box(slide, x, y + SUBTITLE_INSET_Y, w, h)
    for i, line in enumerate(subtitle.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.name, r.font.size = "Arial", Pt(SUBTITLE_FS)

    prs.save(OUT_PPTX)


def main():
    n = read_gene_count()
    subtitle = SUBTITLE.format(n=n)
    with tempfile.TemporaryDirectory() as tmp:
        src = stage(tmp)
        bb = content_bbox(src)
        rect = fit(bb, *CELL)
        doc = render_pdf(src, bb, rect, subtitle)
        render_png(doc)
        render_pptx(src, bb, rect, subtitle, tmp)
    print(f"subtitle gene count: n = {n:,} genes")
    for p in (OUT_PDF, OUT_PNG, OUT_PPTX):
        print(f"wrote {p}  ({os.path.getsize(p)/1e6:.2f} MB)")


if __name__ == "__main__":
    main()
