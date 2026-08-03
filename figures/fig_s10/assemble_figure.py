#!/usr/bin/env python3
"""
Assemble Figure S10 from its individual panels into PDF, PNG and PPTX.

fig_s10 is the higher-in-AD (Direction == 1) counterpart to Fig. 7 panels a-b: the
same dCoPA dot plots, but counting the significant modules whose activity is *higher*
in AD rather than lower. Layout (1 row x 2 cols), MTG on the left, DFC on the right:
  a / b : dCoPA dot plots (panel_{MTG,DFC}_nolegend, from fig_s10.R) + a shared legend

Mirrors fig_7/v8/assemble_figure.py: the layout is computed once; the PDF places panels
as vector (SVG panels converted via rsvg-convert), the PNG is a 300-dpi raster of the
PDF, and the PPTX places 300-dpi panel images at the same positions with live (editable)
title / letters / subtitles. Panel letters, per-panel subtitles and the title are native
text.
"""
import os, subprocess, tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIG  = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_s10")
BASE = os.path.join(FIG, "Kang_Figure_S10_v1")
OUT_PDF, OUT_PNG, OUT_PPTX = BASE + ".pdf", BASE + ".png", BASE + ".pptx"

SRC = {
    "a":      (f"{FIG}/panel_MTG_nolegend.pdf", "pdf"),
    "b":      (f"{FIG}/panel_DFC_nolegend.pdf", "pdf"),
    "legend": (f"{FIG}/panel_legend.svg", "svg"),
}

TITLE = ("Fig. S10 | Up-regulation of modular gene activity in Alzheimer's disease (AD) "
         "occurs mostly in non-neuronal cells.")

# per-panel subtitles (plain a/b, no gold phrase)
SUB = {
    "a": "# of significant dCoPA modules with higher expression in AD MTG",
    "b": "# of significant dCoPA modules with higher expression in AD DFC",
}

# --- page + layout geometry (points; US Letter width, height computed to fit) -------
PW = 612.0
MARG, GUT = 22.0, 12.0
COLW = (PW - 2 * MARG - GUT) / 2.0
XL = MARG
XR = MARG + COLW + GUT

TITLE_TOP, TITLE_H, TITLE_FS = 16.0, 40.0, 10.0
SUB_AB_FS, LET_FS = 6.5, 10.5
HDR_AB = 24.0     # header band height (a/b two-line subtitle)


def content_bbox(src_pdf):
    page = fitz.open(src_pdf)[0]
    r = fitz.Rect(1e9, 1e9, -1e9, -1e9)
    for b in page.get_text("blocks"):
        r |= fitz.Rect(b[:4])
    for d in page.get_drawings():
        r |= d["rect"]
    if r.is_empty or r.x1 < r.x0:
        return page.rect
    return r & page.rect


def fit(bb, cx, cy, cw, ch, halign="left"):
    s = min(cw / bb.width, ch / bb.height)
    dw, dh = bb.width * s, bb.height * s
    dx = cx + (cw - dw) * {"left": 0, "center": .5, "right": 1}[halign]
    return fitz.Rect(dx, cy, dx + dw, cy + dh)


def stage(tmp):
    out = {}
    for k, (path, kind) in SRC.items():
        if kind == "pdf":
            out[k] = path
        else:
            p = os.path.join(tmp, f"{k}.pdf")
            subprocess.run(["rsvg-convert", "-f", "pdf", "-o", p, path], check=True)
            out[k] = p
    return out


def compute_layout(pdf):
    """Return placements shared by every output format, plus the total page height."""
    bb = {k: content_bbox(v) for k, v in pdf.items()}
    panels, letters, subs = [], [], []

    def add(key, x, y, cw, ch, halign="left"):
        panels.append((key, pdf[key], bb[key], fit(bb[key], x, y, cw, ch, halign)))

    # row 1: dot plots + legend sized to the dot-plot scale (matched text size)
    y = TITLE_TOP + TITLE_H + 12
    h1 = COLW / max(bb["a"].width / bb["a"].height, bb["b"].width / bb["b"].height)
    for key, x in (("a", XL), ("b", XR)):
        letters.append((key, x, y + 10)); subs.append((key, x, y, COLW))
        add(key, x, y + HDR_AB, COLW, h1)
    y += HDR_AB + h1

    s_panel = COLW / bb["a"].width                 # scale used for the dot plots
    lw = bb["legend"].width * s_panel              # legend at same scale => same text size
    lh = bb["legend"].height * s_panel
    leg_gap = 11   # vertical gap below the dot plots before the shared legend
    add("legend", (PW - lw) / 2, y + leg_gap, lw, lh, "center")
    y += leg_gap + lh

    ph = y + MARG
    return {"panels": panels, "letters": letters, "subs": subs, "ph": ph}


# ---- PDF (vector) ------------------------------------------------------------
def render_pdf(lay):
    ph = lay["ph"]
    doc = fitz.open()
    page = doc.new_page(width=PW, height=ph)
    page.insert_textbox(fitz.Rect(MARG, TITLE_TOP, PW - MARG, TITLE_TOP + TITLE_H),
                        TITLE, fontname="hebo", fontsize=TITLE_FS,
                        align=fitz.TEXT_ALIGN_LEFT, lineheight=1.15)
    for key, src, bb, rect in lay["panels"]:
        page.show_pdf_page(rect, fitz.open(src), 0, clip=bb)
    for ch, x, y in lay["letters"]:
        page.insert_text((x, y), ch, fontname="hebo", fontsize=LET_FS)
    for key, cx, cy, cw in lay["subs"]:
        # single centered line, vertical middle aligned to the panel letter's middle
        # (letter baseline = cy + 10; a glyph's visual middle ~ baseline - 0.35*fontsize)
        sub = SUB[key]
        w = fitz.get_text_length(sub, "helv", SUB_AB_FS)
        x0 = cx + (cw - w) / 2
        sy = (cy + 10) - 0.35 * (LET_FS - SUB_AB_FS)
        page.insert_text((x0, sy), sub, fontname="helv", fontsize=SUB_AB_FS)
    doc.save(OUT_PDF, deflate=True)
    return doc


def render_png(doc):
    doc[0].get_pixmap(dpi=300).save(OUT_PNG)


# ---- PPTX (panels rasterised at 300 dpi, text editable) ----------------------
def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, bold=False, color=None):
    r = p.add_run(); r.text = text
    r.font.name = "Arial"; r.font.size = Pt(size); r.font.bold = bold
    if color:
        r.font.color.rgb = RGBColor(*color)
    return r


def render_pptx(lay, tmp):
    ph = lay["ph"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(PW), Pt(ph)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    for key, src, bb, rect in lay["panels"]:
        png = os.path.join(tmp, f"{key}_raster.png")
        fitz.open(src)[0].get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72), clip=bb).save(png)
        slide.shapes.add_picture(png, Pt(rect.x0), Pt(rect.y0),
                                 Pt(rect.width), Pt(rect.height))

    tf = _box(slide, MARG, TITLE_TOP, PW - 2 * MARG, TITLE_H)
    _run(tf.paragraphs[0], TITLE, TITLE_FS, bold=True)

    LH = LET_FS + 4
    for ch, x, y in lay["letters"]:
        tf = _box(slide, x, y - LET_FS, LET_FS + 4, LH)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(tf.paragraphs[0], ch, LET_FS, bold=True)

    for key, cx, cy, cw in lay["subs"]:
        tf = _box(slide, cx, cy + 10 - LET_FS, cw, LH)  # same vertical box as the letter
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, SUB[key], SUB_AB_FS)

    prs.save(OUT_PPTX)


def main():
    tmp = tempfile.mkdtemp(prefix="figs10_assemble_")
    pdf = stage(tmp)
    lay = compute_layout(pdf)
    doc = render_pdf(lay)
    render_png(doc)
    render_pptx(lay, tmp)
    for p in (OUT_PDF, OUT_PNG, OUT_PPTX):
        print("Saved:", p)


if __name__ == "__main__":
    main()
