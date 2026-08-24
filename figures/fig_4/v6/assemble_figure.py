#!/usr/bin/env python3
"""
Assemble Figure 4 (v6) from its panels into PDF, PNG and PPTX -- de novo.

Replaces the media-swap assembler, which copied the hand-laid-out fig_4/v5/v5.pptx
and substituted its embedded images. That made a superseded version directory a
build input, and the template was never tracked, so the figure could not be
rebuilt from a clone. It also had to recover which template media held which panel
by content-hashing the embedded SVGs against fig_4/v5, falling back to column
position for five of them -- CELLS below simply says.

Same house pattern as fig_1/v3, fig_3/v4 and fig_6/v6: compute placements once,
emit the PDF as vector, the PNG as a 300-dpi raster of it, and the PPTX with
panels rasterised at 300 dpi and the text live. The placements were read out of
the v6 deliverable the old script produced (v5's boxes plus the aspect-fit and
title respan it applied at build time), then each box was tied to its source file
by md5-matching the images embedded in that deliverable against this folder's
SVGs -- all 19 matched, so the mapping is recovered exactly rather than inferred.

The old path exported the PDF through LibreOffice from a raster-only copy of the
pptx, so every panel arrived rasterised; here the panels stay vector.

Layout (US Letter portrait, 612 x 792 pt): 4 modules x 4 rows, plus 3 summaries.
  a-d  panel_<m>_1.svg                  module expression
  e-h  panel_<m>_2.svg                  bulk correlation
  i-l  panel_<m>_3.svg                  GSEA
  m-p  panel_<m>_4_with_brackets.svg    projection, with the r = ... brackets baked in
  q    panel_Q.svg                      module size distribution
  r    panel_R.svg                      seed-gene mean correlation
  s    panel_S.svg                      REI projection correlation heatmap

The correlation labels are part of the *_with_brackets panels (drawn by
fig_4_v6.R), not slide text, so they track the data rather than needing to be
re-placed by hand. Panel S was a shape picture-fill in the template; here it is
just another placed panel.
"""
import os, subprocess, tempfile

import numpy as np
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt

V6 = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.join(V6, "Kang_Figure_4_v6")
OUT_PDF, OUT_PNG, OUT_PPTX = BASE + ".pdf", BASE + ".png", BASE + ".pptx"

# --- page geometry (points; US Letter portrait) -------------------------------
PW, PH = 612.0, 792.0

TITLE = ("Fig. 4 | Bulk gene coexpression modules in human neocortex reflect "
         "highly reproducible patterns of genomic activity among neocortical "
         "cell types")
TITLE_BOX = (7.2, 5.33, 597.6, 50.4)   # spans the page, wraps to 2 lines
TITLE_FS, LET_FS = 12.0, 11.0

# PowerPoint insets its shape text; the coordinates here are shape origins, so the
# same offsets fig_1/v3, fig_3/v4 and fig_6/v6 needed apply.
INSET_X, INSET_Y_BOX, INSET_Y_LETTER = 7.2, 5.0, 6.25

# Placement cells (x, y, w, h), keyed by source SVG. Several sit at negative x --
# that is the hand layout, not a mistake; those panels bleed off the left margin.
CELLS = {
    "panel_1_1.svg":               (   2.91,  64.80, 143.72,  71.86),
    "panel_2_1.svg":               ( 161.49,  64.80, 143.55,  71.77),
    "panel_3_1.svg":               ( 310.90,  64.80, 143.55,  71.77),
    "panel_4_1.svg":               ( 461.14,  64.80, 143.55,  71.77),
    "panel_1_2.svg":               (  -4.17, 141.11, 143.42,  95.61),
    "panel_2_2.svg":               ( 150.28, 141.17, 143.29,  95.53),
    "panel_3_2.svg":               ( 299.69, 141.17, 143.29,  95.53),
    "panel_4_2.svg":               ( 448.99, 141.17, 143.29,  95.53),
    "panel_1_3.svg":               ( -11.69, 237.54, 161.35,  40.34),
    "panel_2_3.svg":               ( 144.95, 237.71, 160.67,  40.17),
    "panel_3_3.svg":               ( 294.48, 237.60, 160.67,  40.17),
    "panel_4_3.svg":               ( 439.36, 237.60, 160.67,  40.17),
    "panel_1_4_with_brackets.svg": (   2.19, 288.00, 145.37, 110.13),
    "panel_2_4_with_brackets.svg": ( 154.81, 288.00, 145.37, 110.13),
    "panel_3_4_with_brackets.svg": ( 305.06, 288.00, 144.43, 109.42),
    "panel_4_4_with_brackets.svg": ( 454.08, 288.00, 144.43, 109.42),
    "panel_Q.svg":                 (  13.83, 414.00, 138.50,  79.14),
    "panel_R.svg":                 ( 173.58, 414.00, 251.80,  71.94),
    "panel_S.svg":                 ( 441.47, 408.85, 162.39,  93.69),
}

LETTERS = {
    "a": ( -1.50,  45.01), "b": (148.73,  45.01), "c": (296.14,  45.01), "d": (446.37,  45.01),
    "e": ( -1.30, 131.05), "f": (148.93, 131.05), "g": (296.33, 131.05), "h": (446.57, 131.05),
    "i": ( -1.30, 223.68), "j": (148.93, 226.52), "k": (299.17, 226.52), "l": (446.57, 226.52),
    "m": ( -1.30, 273.69), "n": (146.10, 273.69), "o": (299.17, 273.69), "p": (446.57, 273.69),
    "q": ( -1.45, 398.15), "r": (174.44, 398.44), "s": (434.98, 398.15),
}

# (text, x, y, w, size, centred). The module labels are left-aligned over their
# column; the three summary captions are centred over their panels.
CAPTIONS = [
    ("Module 105",  44.65,  52.04,  45.52, 6.0, False),
    ("Module 616", 200.55,  51.96,  45.52, 6.0, False),
    ("Module 627", 350.79,  51.99,  45.52, 6.0, False),
    ("Module 904", 498.19,  52.02,  45.52, 6.0, False),
    ("Distribution of module sizes",                    44.56, 401.75,  93.37, 6.5, True),
    ("Mean pairwise correlations of module seed genes", 230.80, 401.75, 156.78, 6.5, True),
    ("Mean pairwise correlations of REI indices",       452.44, 401.75, 133.00, 6.5, True),
]


def stage(tmp):
    """Panel SVG -> PDF. Inkscape rather than rsvg (see fig_3/v4 and fig_6/v6)."""
    out = {}
    for name in CELLS:
        p = os.path.join(tmp, name.replace(".svg", ".pdf"))
        subprocess.run(["inkscape", os.path.join(V6, name), "--export-type=pdf",
                        f"--export-filename={p}"], check=True, capture_output=True)
        out[name] = p
    return out


def content_bbox(src_pdf):
    """Ink box of a staged panel.

    The text/drawings scan alone is NOT enough here: get_drawings() does not
    descend into Form XObjects, and Inkscape emits panel S's column dendrogram
    inside one. Trimming to that under-reported box clipped the dendrogram off the
    figure. So union it with a raster ink scan, which sees everything the page
    actually paints; the union can only ever grow the box, never crop.
    """
    page = fitz.open(src_pdf)[0]
    r = fitz.Rect(1e9, 1e9, -1e9, -1e9)
    for b in page.get_text("blocks"):
        r |= fitz.Rect(b[:4])
    for d in page.get_drawings():
        r |= d["rect"]

    # Applied as a SAFETY NET, not a co-equal source: the scan rounds outward by up
    # to a pixel and antialiasing bleeds past the true vector edge, so taking it
    # verbatim would nudge every box by a fraction of a point and shift figures that
    # are already published. Only genuinely missed content -- beyond SCAN_TOL -- wins.
    SCAN_DPI, SCAN_TOL = 144, 1.0
    pm = page.get_pixmap(dpi=SCAN_DPI)
    a = np.frombuffer(pm.samples, dtype=np.uint8).reshape(pm.height, pm.width, pm.n)[:, :, :3]
    ys, xs = np.where((a < 250).any(2))
    if len(ys):
        k = 72.0 / SCAN_DPI
        ink = fitz.Rect(page.rect.x0 + xs.min() * k, page.rect.y0 + ys.min() * k,
                        page.rect.x0 + (xs.max() + 1) * k, page.rect.y0 + (ys.max() + 1) * k)
        if ink.x0 < r.x0 - SCAN_TOL: r.x0 = ink.x0
        if ink.y0 < r.y0 - SCAN_TOL: r.y0 = ink.y0
        if ink.x1 > r.x1 + SCAN_TOL: r.x1 = ink.x1
        if ink.y1 > r.y1 + SCAN_TOL: r.y1 = ink.y1

    if r.is_empty or r.x1 < r.x0:
        return page.rect
    return r & page.rect


def fit(bb, cx, cy, cw, ch):
    """Aspect-preserving fit, centred in the cell -- matching the old script."""
    s = min(cw / bb.width, ch / bb.height)
    dw, dh = bb.width * s, bb.height * s
    return fitz.Rect(cx + (cw - dw) / 2, cy + (ch - dh) / 2,
                     cx + (cw - dw) / 2 + dw, cy + (ch - dh) / 2 + dh)


def compute_layout(pdf):
    panels = []
    for name, cell in CELLS.items():
        bb = content_bbox(pdf[name])
        panels.append((name, pdf[name], bb, fit(bb, *cell)))
    return {"panels": panels}


def render_pdf(lay):
    doc = fitz.open()
    page = doc.new_page(width=PW, height=PH)
    x, y, w, h = TITLE_BOX
    page.insert_textbox(fitz.Rect(x + INSET_X, y + INSET_Y_BOX, x + w, y + INSET_Y_BOX + h),
                        TITLE, fontname="hebo", fontsize=TITLE_FS,
                        align=fitz.TEXT_ALIGN_LEFT, lineheight=1.15)
    for _, src, bb, rect in lay["panels"]:
        page.show_pdf_page(rect, fitz.open(src), 0, clip=bb)
    for ch, (lx, ly) in LETTERS.items():
        page.insert_text((lx + INSET_X, ly + INSET_Y_LETTER + LET_FS), ch,
                         fontname="hebo", fontsize=LET_FS)
    for text, cx, cy, cw, fs, centred in CAPTIONS:
        x0 = cx if centred else cx + INSET_X      # centred text centres in the raw box
        page.insert_textbox(fitz.Rect(x0, cy + INSET_Y_BOX, x0 + cw, cy + INSET_Y_BOX + 4 * fs),
                            text, fontname="helv", fontsize=fs, lineheight=1.2,
                            align=fitz.TEXT_ALIGN_CENTER if centred else fitz.TEXT_ALIGN_LEFT)
    doc.save(OUT_PDF, deflate=True)
    return doc


def render_png(doc):
    doc[0].get_pixmap(dpi=300).save(OUT_PNG)


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, bold=False):
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    return r


def render_pptx(lay, tmp):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(PW), Pt(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for name, src, bb, rect in lay["panels"]:
        png = os.path.join(tmp, name.replace(".svg", "_raster.png"))
        fitz.open(src)[0].get_pixmap(matrix=fitz.Matrix(300 / 72, 300 / 72),
                                     clip=bb).save(png)
        slide.shapes.add_picture(png, Pt(rect.x0), Pt(rect.y0),
                                 Pt(rect.width), Pt(rect.height))

    x, y, w, h = TITLE_BOX
    _run(_box(slide, x + INSET_X, y + INSET_Y_BOX, w, h).paragraphs[0],
         TITLE, TITLE_FS, bold=True)
    for ch, (lx, ly) in LETTERS.items():
        _run(_box(slide, lx + INSET_X, ly + INSET_Y_BOX, 22.0,
                  LET_FS * 1.7).paragraphs[0], ch, LET_FS, bold=True)
    for text, cx, cy, cw, fs, centred in CAPTIONS:
        _run(_box(slide, cx if centred else cx + INSET_X, cy + INSET_Y_BOX, cw,
                  4 * fs).paragraphs[0], text, fs)

    prs.save(OUT_PPTX)


def main():
    with tempfile.TemporaryDirectory() as tmp:
        pdf = stage(tmp)
        lay = compute_layout(pdf)
        doc = render_pdf(lay)
        render_png(doc)
        render_pptx(lay, tmp)
    for p in (OUT_PDF, OUT_PNG, OUT_PPTX):
        print(f"wrote {p}  ({os.path.getsize(p)/1e6:.2f} MB)")


if __name__ == "__main__":
    main()
