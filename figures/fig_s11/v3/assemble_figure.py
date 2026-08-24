#!/usr/bin/env python3
"""
Figure assembly for Fig. S11 (v3 = post-HGNC-fix re-run of v2; panels regenerated
from the updated fig_7/v7.1 dCoPA gene lists. Layout is read from v2/v2.pptx, which
is unchanged; only the panel content differs).

Deliverables (all named Kang_Figure_S11_v3.*):
  .pdf   — vector
  .png   — raster, 300 dpi
  .pptx  — one slide, the figure embedded as a vector SVG with a PNG fallback

Layout: PANELS and LABELS below ARE the layout. They were read once out of the
hand-laid-out fig_s11/v2/v2.pptx and frozen here, so a superseded version folder
is no longer a build input and the figure rebuilds from a clone (the template was
never tracked). This script composes a single master SVG that INLINES the panel
SVGs (so they stay vector)
with the title / labels / panel titles drawn as Arial <text>. That master SVG is
an internal intermediate (deleted at the end, not a deliverable): rsvg-convert
renders it to the PDF and PNG, and it is embedded as the pptx's vector layer.

Panels (produced by make_panels.R in this folder), MTG left / DFC right:
  a_MTG  -> panel_a_MTG.svg   (Fig. S11a, left)
  b_DFC  -> panel_b_DFC.svg   (Fig. S11b, right)
  legend.svg                  (shared -log10(FDR) legend, transparent, drawn twice)

The legend is drawn ONCE, transparent, centred in the gap between the two panels
so it never occludes a heatmap (the v2 reference overlapped each panel with an
opaque legend; that is the bug this fixes).
"""
import os, re, subprocess

V3       = os.path.join(os.environ.get("REPO_DIR", "/home/gugene/code/git/kang-oldham-2026"), "figures/fig_s11/v3")  # panels + deliverables
OUT_PDF  = f"{V3}/Kang_Figure_S11_v3.pdf"
OUT_PNG  = f"{V3}/Kang_Figure_S11_v3.png"
OUT_PPTX = f"{V3}/Kang_Figure_S11_v3.pptx"
# composed master SVG — internal intermediate only (drives the vector PDF and the
# pptx's vector layer); deleted at the end so it is not a deliverable.
MASTER_SVG = f"{V3}/.s11_master.svg"

# panel key -> source SVG in this folder (left=MTG, right=DFC)
MEDIA_SRC = {"a_MTG": "panel_a_MTG.svg",
             "b_DFC": "panel_b_DFC.svg"}
LEGEND_SRC = "legend.svg"

TITLE = ("Fig. S11 | Modular gene activity is reproducibly and coordinately "
         "down-regulated in specific neuronal subclasses.")

EMU_PER_PX = 9525.0          # 96 dpi
def px(emu): return float(emu) / EMU_PER_PX

# ---- layout, frozen from v2.pptx (EMU; 914400 per inch) ------------------
# Panel boxes, and the text shapes drawn over them. Sizes are px (the pptx sz/100
# converted at 96/72). Everything else on the canvas is derived from these.
pics = {
    "a_MTG": (  82440, 887040, 3200040, 3008160),
    "b_DFC": (4003920, 887040, 3200040, 3008160),
}
# (text, x, y, cx, cy, size_px, bold, centred)
texts = [
    ("Significance of overlap of dCoPA genes (MTG)", 1135080, 665640, 2006640, 196200, 9.3333, False, True),
    ("a",                                             468000, 608400,  279000, 303120, 18.6667, True, False),
    ("b",                                            4140000, 608400,  289080, 303120, 18.6667, True, False),
    ("Significance of overlap of dCoPA genes (DFC)", 4852080, 665640, 1990800, 196200, 9.3333, False, True),
]

# Legends: ONE per panel (a and b), transparent. The legend's visible colour-bar
# block is ~0.435x its box width, centred in the box. Place panel a's legend so that
# block is centred in the gap between the panels, and panel b's in a right-hand
# margin, so neither colour-bar overlaps a heatmap.
panel_media = ["a_MTG", "b_DFC"]
leg_w, leg_h = 1218240, 913320
leg_y = (pics["a_MTG"][1] + pics["a_MTG"][3] / 2) - leg_h / 2
LEG_CONTENT_HALF = 0.435 / 2 * leg_w     # half-width of the visible colour-bar block
GAP_CLEAR = 95250                        # 0.1 in clearance from a panel edge
a_r = pics["a_MTG"][0] + pics["a_MTG"][2]   # panel a right edge
b_l = pics["b_DFC"][0]                            # panel b left edge
b_r = pics["b_DFC"][0] + pics["b_DFC"][2]   # panel b right edge
legA_center = (a_r + b_l) / 2                          # colour-bar centred in the gap
legB_center = b_r + GAP_CLEAR + LEG_CONTENT_HALF       # colour-bar clear, right of panel b
legends = [("legendA_", legA_center - leg_w / 2),      # panel a (MTG)
           ("legendB_", legB_center - leg_w / 2)]      # panel b (DFC)

# ---- canvas: crop to content (down to the panel bottoms; wide enough for legend B)
MARGIN = 91440  # 0.1 in
right  = max(legB_center + LEG_CONTENT_HALF,
             max(pics[m][0] + pics[m][2] for m in panel_media))
bottom = max(pics[m][1] + pics[m][3] for m in panel_media)
W_EMU, H_EMU = right + MARGIN, bottom + MARGIN
W, H = px(W_EMU), px(H_EMU)

# ---- SVG inlining (keeps panels vector; prefix ids to avoid cross-panel clashes)
def _idprefix(s, p):
    s = re.sub(r'\bid="([^"]+)"',          lambda m: f'id="{p}{m.group(1)}"', s)
    s = re.sub(r'xlink:href="#([^"]+)"',   lambda m: f'xlink:href="#{p}{m.group(1)}"', s)
    s = re.sub(r'(?<!xlink:)\bhref="#([^"]+)"', lambda m: f'href="#{p}{m.group(1)}"', s)
    s = re.sub(r'url\(#([^)]+)\)',         lambda m: f'url(#{p}{m.group(1)})', s)
    return s

def inline(path, prefix, x, y, w, h):
    s = open(path, encoding="utf-8").read()
    s = re.sub(r'<\?xml[^>]*\?>', '', s)
    s = re.sub(r'<!DOCTYPE[^>]*>', '', s).strip()
    vb = re.search(r'viewBox=["\']([^"\']+)["\']', s).group(1)
    s = _idprefix(s, prefix)
    s = re.sub(r'<svg\b[^>]*>',
               f'<svg x="{x:.2f}" y="{y:.2f}" width="{w:.2f}" height="{h:.2f}" '
               f'viewBox="{vb}" preserveAspectRatio="xMidYMid meet" overflow="visible">',
               s, count=1)
    return s

def esc(t): return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
FAMILY = "Arial, 'Liberation Sans', sans-serif"

parts = [f'<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" '
         f'width="{W:.2f}" height="{H:.2f}" viewBox="0 0 {W:.2f} {H:.2f}">',
         f'<rect width="{W:.2f}" height="{H:.2f}" fill="white"/>']

# panels first (opaque), then the transparent legends on top (one per panel)
for media in panel_media:
    x, y, w, h = pics[media]
    parts.append(inline(f"{V3}/{MEDIA_SRC[media]}", media + "_",
                        px(x), px(y), px(w), px(h)))
for prefix, lx in legends:
    parts.append(inline(f"{V3}/{LEGEND_SRC}", prefix,
                        px(lx), px(leg_y), px(leg_w), px(leg_h)))

# title: left-justified, flush with the left edge of panel a, across the top
t_size = 14.0
parts.append(f'<text x="{px(pics["a_MTG"][0]):.1f}" y="{px(31680)+t_size*0.9:.1f}" text-anchor="start" '
             f'font-family="{FAMILY}" font-size="{t_size:.1f}" font-weight="bold">{esc(TITLE)}</text>')

# panel letters (a, b) and per-panel titles, from the v2 text shapes, in Arial.
# Centred subtitles are centred over the heatmap GRID (the colour-tile box, excluding
# the axis text), read straight from each panel SVG's plot-area rect.
def _grid_center_px(panel_media):
    s = open(f"{V3}/{MEDIA_SRC[panel_media]}", encoding="utf-8").read()
    m = re.search(r"viewBox=['\"][\d.]+ [\d.]+ ([\d.]+) ([\d.]+)", s)
    vbw, vbh = float(m.group(1)), float(m.group(2))
    rects = [(float(rx), float(rw))
             for rx, rw in re.findall(r"<rect x='([\d.]+)'[^>]*width='([\d.]+)'", s)]
    gx, gw = max((r for r in rects if r[1] < 0.9 * vbw), key=lambda r: r[1])  # plot area
    bx, by, bw, bh = pics[panel_media]
    scale = min(bw / vbw, bh / vbh)                 # xMidYMid meet
    offx = bx + (bw - vbw * scale) / 2
    return px(offx + (gx + gw / 2) * scale)
GRID_CX = {"(MTG)": _grid_center_px("a_MTG"), "(DFC)": _grid_center_px("b_DFC")}
for txt, x, y, cx, cy, sz, bold, centred in texts:
    weight = ' font-weight="bold"' if bold else ''
    if centred:
        key = next((k for k in GRID_CX if k in txt), None)
        tx, anchor = (GRID_CX[key] if key else px(x + cx / 2)), "middle"
    else:
        tx, anchor = px(x) + 2, "start"
    ty = px(y) + sz * 0.95   # box-top -> approx baseline (inset + ascent)
    parts.append(f'<text x="{tx:.1f}" y="{ty:.1f}" text-anchor="{anchor}" '
                 f'font-family="{FAMILY}" font-size="{sz:.1f}"{weight}>{esc(txt)}</text>')

parts.append("</svg>")
open(MASTER_SVG, "w", encoding="utf-8").write("\n".join(parts))   # intermediate only

# vector PDF (deliverable)
subprocess.run(["rsvg-convert", "-f", "pdf", MASTER_SVG, "-o", OUT_PDF], check=True)
print("wrote:", OUT_PDF, f"({W:.0f}x{H:.0f}px)")

# raster PNG (deliverable). Also serves as the pptx's fallback blip below.
# The master SVG is sized in px at 96 dpi, so rsvg's --dpi is a no-op; scale the
# pixel dimensions instead to render at 300 dpi.
PNG_DPI = 300
png_w, png_h = round(W * PNG_DPI / 96), round(H * PNG_DPI / 96)
subprocess.run(["rsvg-convert", "-f", "png", "-w", str(png_w), "-h", str(png_h),
                MASTER_SVG, "-o", OUT_PNG], check=True)
print("wrote:", OUT_PNG, f"({png_w}x{png_h}px @ {PNG_DPI}dpi)")

# ---- PPTX (deliverable): one slide = the composed figure, embedded as a vector
# SVG with a PNG fallback (the native PowerPoint pattern: a:blip -> PNG,
# asvg:svgBlip -> SVG), so the slide stays vector in PowerPoint while remaining
# openable everywhere. The SVG lives only inside the pptx; no .svg file is emitted.
from pptx import Presentation
from pptx.util import Emu
from pptx.opc.package import Part
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(int(W_EMU)), Emu(int(H_EMU))
slide = prs.slides.add_slide(prs.slide_layouts[6])          # blank layout
pic = slide.shapes.add_picture(OUT_PNG, 0, 0, Emu(int(W_EMU)), Emu(int(H_EMU)))

# add the master SVG as a media part related from the slide, then point an
# asvg:svgBlip at it inside the picture's a:blip (alongside the PNG fallback).
spart = slide.part
svg_partname = spart.package.next_partname("/ppt/media/image%d.svg")
svg_part = Part(svg_partname, "image/svg+xml",
                spart.package, open(MASTER_SVG, "rb").read())
IMAGE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
svg_rid = spart.relate_to(svg_part, IMAGE_REL)

blip = pic._element.blipFill.find(qn("a:blip"))
blip.append(parse_xml(
    '<a:extLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">'
    '<asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    f'r:embed="{svg_rid}"/></a:ext></a:extLst>'))

prs.save(OUT_PPTX)
os.remove(MASTER_SVG)               # intermediate; not a deliverable
print("wrote:", OUT_PPTX)
